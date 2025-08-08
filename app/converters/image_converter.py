"""
图片转换器
"""
import asyncio
import tempfile
import os
from typing import Dict, Any, Optional
from pathlib import Path
from loguru import logger

try:
    from PIL import Image
    import pytesseract
    from docx import Document
    from docx.shared import Inches, Pt
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from openpyxl import Workbook
    from openpyxl.drawing.image import Image as XLImage
    from pptx import Presentation
    from pptx.util import Inches as PptInches
    from pptx.enum.shapes import MSO_SHAPE
except ImportError as e:
    logger.warning(f"图片转换器依赖库未安装: {e}")


class ImageConverter:
    """图片转换器类"""
    
    def __init__(self):
        """初始化图片转换器"""
        self.supported_formats = ['pdf', 'docx', 'xlsx', 'pptx', 'jpg', 'png', 'gif', 'bmp', 'tiff']
    
    async def image_to_pdf(
        self, 
        input_path: str, 
        output_path: str, 
        options: Optional[Dict[str, Any]] = None
    ) -> Dict[str, Any]:
        """
        图片转换为 PDF
        
        Args:
            input_path: 输入文件路径
            output_path: 输出文件路径
            options: 转换选项
            
        Returns:
            转换结果
        """
        try:
            # 打开图片
            with Image.open(input_path) as image:
                # 转换为 RGB 模式
                if image.mode != 'RGB':
                    image = image.convert('RGB')
                
                # 保存为 PDF
                image.save(output_path, 'PDF', resolution=300.0)
                
                return {
                    'success': True,
                    'message': '图片转 PDF 成功'
                }
        except Exception as e:
            logger.error(f"图片转 PDF 失败: {e}")
            return {
                'success': False,
                'error': f'图片转 PDF 失败: {str(e)}'
            }
    
    async def image_to_office(
        self, 
        input_path: str, 
        output_path: str, 
        target_format: str, 
        options: Optional[Dict[str, Any]] = None
    ) -> Dict[str, Any]:
        """
        图片转换为 Office 文档
        
        Args:
            input_path: 输入文件路径
            output_path: 输出文件路径
            target_format: 目标格式
            options: 转换选项
            
        Returns:
            转换结果
        """
        try:
            # 根据目标格式选择转换方法
            if target_format == 'docx':
                return await self._image_to_docx(input_path, output_path, options)
            elif target_format == 'pptx':
                return await self._image_to_pptx(input_path, output_path, options)
            elif target_format == 'xlsx':
                return await self._image_to_xlsx(input_path, output_path, options)
            else:
                return {
                    'success': False,
                    'error': f'不支持的 Office 格式: {target_format}'
                }
        except Exception as e:
            logger.error(f"图片转 Office 失败: {e}")
            return {
                'success': False,
                'error': f'图片转 Office 失败: {str(e)}'
            }
    
    async def _image_to_docx(
        self, 
        input_path: str, 
        output_path: str, 
        options: Optional[Dict[str, Any]] = None
    ) -> Dict[str, Any]:
        """图片转换为 Word 文档"""
        try:
            # 创建新的 Word 文档
            doc = Document()
            
            # 设置页面边距为最小
            sections = doc.sections
            for section in sections:
                section.top_margin = Inches(0.5)
                section.bottom_margin = Inches(0.5)
                section.left_margin = Inches(0.5)
                section.right_margin = Inches(0.5)
            
            # 只添加图片，不添加任何文字
            paragraph = doc.add_paragraph()
            run = paragraph.add_run()
            
            # 获取图片尺寸
            with Image.open(input_path) as img:
                width, height = img.size
                
                # 计算合适的显示尺寸（最大宽度 7 英寸，适应页面）
                max_width = 7.0
                if width > height:
                    # 横向图片
                    display_width = min(max_width, width / 100)
                    display_height = (height / width) * display_width
                else:
                    # 纵向图片
                    display_height = min(max_width * 1.2, height / 100)
                    display_width = (width / height) * display_height
            
            # 插入图片
            run.add_picture(input_path, width=Inches(display_width))
            
            # 保存文档
            doc.save(output_path)
            
            return {
                'success': True,
                'message': '图片转 Word 文档成功',
                'image_size': f'{width}x{height}',
                'output_format': 'docx'
            }
            
        except Exception as e:
            logger.error(f"图片转 Word 失败: {e}")
            return {
                'success': False,
                'error': f'图片转 Word 失败: {str(e)}'
            }
    
    async def _image_to_pptx(
        self, 
        input_path: str, 
        output_path: str, 
        options: Optional[Dict[str, Any]] = None
    ) -> Dict[str, Any]:
        """图片转换为 PowerPoint 演示文稿"""
        try:
            # 创建新的 PowerPoint 演示文稿
            prs = Presentation()
            
            # 使用空白页布局，只添加图片
            blank_slide_layout = prs.slide_layouts[6]  # 空白页布局
            slide = prs.slides.add_slide(blank_slide_layout)
            
            # 获取图片尺寸
            with Image.open(input_path) as img:
                width, height = img.size
                
                # 计算合适的显示尺寸（最大宽度 9 英寸，充分利用幻灯片空间）
                max_width = 9.0
                max_height = 6.5
                
                if width > height:
                    # 横向图片
                    display_width = min(max_width, width / 100)
                    display_height = (height / width) * display_width
                    if display_height > max_height:
                        display_height = max_height
                        display_width = (width / height) * display_height
                else:
                    # 纵向图片
                    display_height = min(max_height, height / 100)
                    display_width = (width / height) * display_height
                    if display_width > max_width:
                        display_width = max_width
                        display_height = (height / width) * display_width
            
            # 计算图片位置（居中）
            slide_width = 10.0  # 幻灯片宽度（英寸）
            slide_height = 7.5   # 幻灯片高度（英寸）
            
            left = (slide_width - display_width) / 2
            top = (slide_height - display_height) / 2
            
            # 添加图片
            slide.shapes.add_picture(
                input_path, 
                left=PptInches(left), 
                top=PptInches(top),
                width=PptInches(display_width),
                height=PptInches(display_height)
            )
            
            # 保存演示文稿
            prs.save(output_path)
            
            return {
                'success': True,
                'message': '图片转 PowerPoint 演示文稿成功',
                'image_size': f'{width}x{height}',
                'output_format': 'pptx',
                'slides_count': len(prs.slides)
            }
            
        except Exception as e:
            logger.error(f"图片转 PowerPoint 失败: {e}")
            return {
                'success': False,
                'error': f'图片转 PowerPoint 失败: {str(e)}'
            }
    
    async def _image_to_xlsx(
        self, 
        input_path: str, 
        output_path: str, 
        options: Optional[Dict[str, Any]] = None
    ) -> Dict[str, Any]:
        """图片转换为 Excel 工作簿"""
        try:
            # 创建新的 Excel 工作簿
            wb = Workbook()
            
            # 获取默认工作表，只用来放图片
            ws = wb.active
            ws.title = "图片"
            
            # 获取图片尺寸
            with Image.open(input_path) as img:
                width, height = img.size
            
            # 添加图片到工作表
            try:
                # 直接使用原始图片文件
                img_excel = XLImage(input_path)
                
                # 设置合适的显示尺寸，充分利用 Excel 空间
                if width > height:
                    img_excel.width = 500
                    img_excel.height = int((height / width) * 500)
                else:
                    img_excel.height = 400
                    img_excel.width = int((width / height) * 400)
                
                # 限制最大尺寸
                if img_excel.width > 600:
                    img_excel.width = 600
                if img_excel.height > 500:
                    img_excel.height = 500
                
                # 将图片放在 A1 位置
                ws.add_image(img_excel, 'A1')
                    
            except Exception as e:
                logger.warning(f"添加图片到 Excel 失败: {e}")
                ws['A1'] = "图片加载失败"
            
            # 保存工作簿
            wb.save(output_path)
            
            return {
                'success': True,
                'message': '图片转 Excel 工作簿成功',
                'image_size': f'{width}x{height}',
                'output_format': 'xlsx',
                'sheets_count': len(wb.sheetnames)
            }
            
        except Exception as e:
            logger.error(f"图片转 Excel 失败: {e}")
            return {
                'success': False,
                'error': f'图片转 Excel 失败: {str(e)}'
            }
    
    async def image_to_image(
        self, 
        input_path: str, 
        output_path: str, 
        target_format: str, 
        options: Optional[Dict[str, Any]] = None
    ) -> Dict[str, Any]:
        """
        图片格式间转换
        
        Args:
            input_path: 输入文件路径
            output_path: 输出文件路径
            target_format: 目标格式
            options: 转换选项
            
        Returns:
            转换结果
        """
        try:
            # 打开图片
            with Image.open(input_path) as image:
                # 处理图片模式
                if image.mode in ('RGBA', 'LA'):
                    # 转换为 RGB
                    background = Image.new('RGB', image.size, (255, 255, 255))
                    background.paste(image, mask=image.split()[-1] if image.mode == 'RGBA' else None)
                    image = background
                elif image.mode != 'RGB':
                    image = image.convert('RGB')
                
                # 设置图片质量
                quality = options.get('quality', 'medium') if options else 'medium'
                if quality == 'high':
                    save_quality = 95
                elif quality == 'low':
                    save_quality = 60
                else:
                    save_quality = 80
                
                # 调整图片尺寸
                if options and options.get('output_size') and options.get('output_size') != 'original':
                    image = self._resize_image(image, options['output_size'])
                
                # 保存图片
                # 处理格式映射
                save_format = target_format.upper()
                if save_format == 'JPG':
                    save_format = 'JPEG'
                elif save_format == 'TIFF':
                    save_format = 'TIFF'
                
                # 设置保存参数
                save_kwargs = {}
                if save_format == 'JPEG':
                    save_kwargs['quality'] = save_quality
                    save_kwargs['optimize'] = True
                elif save_format == 'PNG':
                    save_kwargs['optimize'] = True
                elif save_format == 'GIF':
                    save_kwargs['optimize'] = True
                
                image.save(output_path, save_format, **save_kwargs)
                
                return {
                    'success': True,
                    'message': f'图片转 {target_format.upper()} 成功'
                }
                
        except Exception as e:
            logger.error(f"图片格式转换失败: {e}")
            return {
                'success': False,
                'error': f'图片格式转换失败: {str(e)}'
            }
    
    async def pdf_to_image(
        self, 
        input_path: str, 
        output_path: str, 
        target_format: str, 
        options: Optional[Dict[str, Any]] = None
    ) -> Dict[str, Any]:
        """
        PDF 转换为图片
        
        Args:
            input_path: 输入文件路径
            output_path: 输出文件路径
            target_format: 目标格式
            options: 转换选项
            
        Returns:
            转换结果
        """
        try:
            # 这里调用 PDF 转换器的方法
            from .pdf_converter import PDFConverter
            pdf_converter = PDFConverter()
            return await pdf_converter.pdf_to_image(input_path, output_path, target_format, options)
        except Exception as e:
            logger.error(f"PDF 转图片失败: {e}")
            return {
                'success': False,
                'error': f'PDF 转图片失败: {str(e)}'
            }
    
    async def office_to_image(
        self, 
        input_path: str, 
        output_path: str, 
        target_format: str, 
        options: Optional[Dict[str, Any]] = None
    ) -> Dict[str, Any]:
        """
        Office 文档转换为图片
        
        Args:
            input_path: 输入文件路径
            output_path: 输出文件路径
            target_format: 目标格式
            options: 转换选项
            
        Returns:
            转换结果
        """
        try:
            # TODO: 实现 Office 转图片功能
            return {
                'success': False,
                'error': 'Office 转图片功能暂未实现'
            }
        except Exception as e:
            logger.error(f"Office 转图片失败: {e}")
            return {
                'success': False,
                'error': f'Office 转图片失败: {str(e)}'
            }
    
    def _resize_image(self, image: Image.Image, output_size: str) -> Image.Image:
        """调整图片尺寸"""
        if output_size == 'A4':
            # A4 尺寸: 210mm x 297mm, 300 DPI
            target_width = int(210 * 300 / 25.4)  # mm to pixels
            target_height = int(297 * 300 / 25.4)
        elif output_size == 'letter':
            # Letter 尺寸: 8.5" x 11", 300 DPI
            target_width = int(8.5 * 300)
            target_height = int(11 * 300)
        else:
            return image
        
        # 保持宽高比
        image.thumbnail((target_width, target_height), Image.Resampling.LANCZOS)
        return image
    
    def get_image_info(self, input_path: str) -> Dict[str, Any]:
        """获取图片信息"""
        try:
            with Image.open(input_path) as image:
                return {
                    'format': image.format,
                    'mode': image.mode,
                    'size': image.size,
                    'width': image.width,
                    'height': image.height,
                    'file_size': image.size[0] * image.size[1] * len(image.getbands())
                }
        except Exception as e:
            logger.error(f"获取图片信息失败: {e}")
            return {} 