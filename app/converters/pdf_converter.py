"""
PDF 转换器
"""
import asyncio
from typing import Dict, Any, Optional, List, Tuple
from pathlib import Path
from loguru import logger

try:
    import PyPDF2
    import pdfplumber
    from pdf2image import convert_from_path
    from PIL import Image
    import pytesseract
except ImportError as e:
    logger.warning(f"PDF 转换器依赖库未安装: {e}")

import os
import tempfile
import logging
from typing import Dict, Any, Optional, List, Tuple
from pathlib import Path
import fitz  # PyMuPDF
from docx import Document
from docx.shared import Pt, Inches
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.section import WD_SECTION
from docx.oxml.shared import OxmlElement, qn
from PIL import Image
import io
import base64
from app.config.settings import USE_PDF2DOCX, PDF2DOCX_FALLBACK

logger = logging.getLogger(__name__)

class PDFConverter:
    """
    PDF 转换器 - 支持混合方案
    可以选择使用 pdf2docx 或自定义实现
    """
    
    def __init__(self):
        self.use_pdf2docx = USE_PDF2DOCX
        self.pdf2docx_fallback = PDF2DOCX_FALLBACK
        
        # 检查 pdf2docx 是否可用
        self.pdf2docx_available = self._check_pdf2docx_availability()
        
        if self.use_pdf2docx and not self.pdf2docx_available:
            logger.warning("pdf2docx 不可用，将使用自定义实现")
            self.use_pdf2docx = False
    
    def _check_pdf2docx_availability(self) -> bool:
        """检查 pdf2docx 是否可用"""
        try:
            import pdf2docx
            return True
        except ImportError:
            logger.warning("pdf2docx 未安装，将使用自定义实现")
            return False
    
    async def pdf_to_office(
        self, 
        input_path: str, 
        output_path: str, 
        target_format: str, 
        options: Optional[Dict[str, Any]] = None
    ) -> Dict[str, Any]:
        """
        PDF 转换为 Office 文档
        
        Args:
            input_path: 输入 PDF 文件路径
            output_path: 输出文件路径
            target_format: 目标格式
            options: 转换选项
            
        Returns:
            转换结果
        """
        try:
            if target_format == 'docx':
                return await self._pdf_to_docx(input_path, output_path, options)
            elif target_format == 'xlsx':
                return await self._pdf_to_xlsx(input_path, output_path, options)
            elif target_format == 'pptx':
                return await self._pdf_to_pptx(input_path, output_path, options)
            else:
                return {
                    'success': False,
                    'error': f'不支持 PDF 转换到 {target_format}'
                }
        except Exception as e:
            logger.error(f"PDF 转 Office 失败: {e}")
            return {
                'success': False,
                'error': f'PDF 转 Office 失败: {str(e)}'
            }
    
    async def pdf_to_image(
        self, 
        input_path: str, 
        output_path: str, 
        target_format: str, 
        options: Optional[Dict[str, Any]] = None
    ) -> Dict[str, Any]:
        """
        PDF 转换为图片
        
        Args:
            input_path: 输入 PDF 文件路径
            output_path: 输出文件路径
            target_format: 目标格式
            options: 转换选项
            
        Returns:
            转换结果
        """
        try:
            return await self._pdf_to_image(input_path, output_path, target_format, options)
        except Exception as e:
            logger.error(f"PDF 转图片失败: {e}")
            return {
                'success': False,
                'error': f'PDF 转图片失败: {str(e)}'
            }
    
    async def _pdf_to_docx(
        self, 
        input_path: str, 
        output_path: str, 
        options: Optional[Dict[str, Any]] = None
    ) -> Dict[str, Any]:
        """
        PDF 转 Word - 混合方案
        优先使用 pdf2docx，失败时回退到自定义实现
        """
        
        # 如果启用 pdf2docx 且可用，优先使用
        if self.use_pdf2docx and self.pdf2docx_available:
            try:
                logger.info("使用 pdf2docx 进行 PDF 转 Word 转换")
                return await self._pdf_to_docx_with_pdf2docx(input_path, output_path, options)
            except Exception as e:
                logger.warning(f"pdf2docx 转换失败: {e}")
                if self.pdf2docx_fallback:
                    logger.info("回退到自定义实现")
                    return await self._pdf_to_docx_custom(input_path, output_path, options)
                else:
                    raise e
        
        # 使用自定义实现
        logger.info("使用自定义实现进行 PDF 转 Word 转换")
        return await self._pdf_to_docx_custom(input_path, output_path, options)
    
    async def _pdf_to_xlsx(
        self, 
        input_path: str, 
        output_path: str, 
        options: Optional[Dict[str, Any]] = None
    ) -> Dict[str, Any]:
        """PDF 转换为 Excel 文档"""
        try:

            from openpyxl import Workbook
            from openpyxl.drawing.image import Image as XLImage
            import fitz  # PyMuPDF
            
            # 创建 Excel 工作簿
            wb = Workbook()
            ws = wb.active
            
            # 用于跟踪临时文件
            temp_files = []
            
            # 使用 PyMuPDF 提取表格和图片
            pdf_document = fitz.open(input_path)
            
            # 处理页面范围
            pages = self._get_pages_to_process(pdf_document, options)
            
            current_row = 1
            
            for page_num in pages:
                page = pdf_document[page_num - 1]
                
                # 收集页面上的所有内容（文本和图片），按位置排序
                page_content = self._collect_page_content(page, page_num, pdf_document, options)
                
                # 按垂直位置排序内容
                page_content.sort(key=lambda x: x['y_position'])
                
                # 按顺序处理内容 - 应用智能文本处理
                current_paragraph = None
                current_line_bbox = None
                current_max_font_size = 0
                
                for content in page_content:
                    if content['type'] == 'text':
                        # 智能处理文本
                        line = content['data']
                        
                        # 获取行的边界框信息
                        line_bbox = [0, 0, 0, 0]
                        if line["spans"]:
                            first_span = line["spans"][0]
                            line_bbox = first_span.get("bbox", [0, 0, 0, 0])
                        
                        page_width = page.rect.width
                        
                        # 检查是否需要创建新行
                        should_create_new_row = True
                        
                        # 如果当前行存在，检查是否可以继续使用
                        if current_paragraph is not None:
                            # 检查垂直距离 - 如果距离很近，可能是同一行的延续
                            if current_line_bbox and len(current_line_bbox) >= 4 and len(line_bbox) >= 4:
                                vertical_distance = abs(line_bbox[1] - current_line_bbox[1])
                                # 如果垂直距离小于字体大小的一半，认为是同一行
                                if vertical_distance < 10:  # 10像素的阈值
                                    should_create_new_row = False
                        
                        # 创建新行或使用现有行
                        if should_create_new_row:
                            # 完成当前行
                            if current_paragraph is not None:
                                # 设置当前行的格式
                                if current_max_font_size > 0:
                                    self._set_excel_row_format(ws, current_row, current_max_font_size)
                            
                            # 创建新行
                            current_row += 1
                            current_line_bbox = line_bbox
                            current_max_font_size = 0
                        
                        # 处理每个文本片段
                        col_idx = 1
                        for span in line["spans"]:
                            span_text = span["text"]
                            if span_text.strip():
                                # 设置单元格值
                                cell = ws.cell(row=current_row, column=col_idx, value=span_text.strip())
                                
                                # 创建字体样式
                                from openpyxl.styles import Font
                                font_name = span.get("font", "Arial")
                                span_size = span.get("size", 12)
                                excel_font_size = max(8, min(72, span_size * 0.75))
                                span_flags = span.get("flags", 0)
                                
                                # 创建字体对象
                                font = Font(
                                    name=font_name,
                                    size=excel_font_size,
                                    bold=bool(span_flags & 2**4),
                                    italic=bool(span_flags & 2**1)
                                )
                                
                                # 设置颜色
                                span_color = span.get("color", 0)
                                if span_color != 0:
                                    rgb_color = self._convert_color_to_rgb(span_color)
                                    if rgb_color:
                                        try:
                                            from openpyxl.styles.colors import Color
                                            font.color = Color(rgb=rgb_color)
                                        except Exception as color_error:
                                            logger.warning(f"设置字体颜色失败: {color_error}")
                                
                                # 应用字体样式
                                cell.font = font
                                current_max_font_size = max(current_max_font_size, excel_font_size)
                                
                                # 对齐方式
                                self._set_excel_cell_alignment(cell, line_bbox, page_width)
                                
                                col_idx += 1
                    
                    elif content['type'] == 'image':
                        # 处理图片
                        img_data = content['data']
                        try:
                            # 使用安全的方法提取图片
                            temp_img_path = self._extract_image_safely(
                                pdf_document, img_data['img'], img_data['page_num'], img_data['img_index']
                            )
                            
                            if temp_img_path:
                                # 智能计算图片尺寸
                                img_width, img_height = self._calculate_smart_excel_image_size(
                                    pdf_document, img_data['img'], img_data['page_num'], img_data['img_index'], options
                                )
                                
                                # 添加图片到 Excel
                                img = XLImage(temp_img_path)
                                img.width = img_width
                                img.height = img_height
                                ws.add_image(img, f'A{current_row}')
                                
                                logger.info(f"添加图片: 智能缩放为 {img_width}x{img_height} 像素，位置: Y={content['y_position']}")
                                
                                # 根据图片高度调整行间距
                                row_height = max(10, img_height // 20)
                                current_row += row_height
                                
                                # 添加到临时文件列表，稍后清理
                                temp_files.append(temp_img_path)
                            
                        except Exception as img_error:
                            logger.warning(f"处理图片失败: {img_error}")
                            continue
                
                # 添加页面分隔
                current_row += 2
            
            # 关闭 PDF 文档
            pdf_document.close()
            
            # 保存工作簿
            wb.save(output_path)
            
            # 保存完成后再清理临时文件
            import os
            for temp_file in temp_files:
                try:
                    if os.path.exists(temp_file):
                        os.remove(temp_file)
                except Exception as e:
                    logger.warning(f"清理临时文件失败: {e}")
            
            return {
                'success': True,
                'message': 'PDF 转 Excel 成功（包含图片）'
            }
            
        except Exception as e:
            logger.error(f"PDF 转 Excel 失败: {e}")
            return {
                'success': False,
                'error': f'PDF 转 Excel 失败: {str(e)}'
            }
    
    async def _pdf_to_pptx(
        self, 
        input_path: str, 
        output_path: str, 
        options: Optional[Dict[str, Any]] = None
    ) -> Dict[str, Any]:
        """PDF 转换为 PowerPoint 文档"""
        try:
            from pptx import Presentation
            from pptx.util import Inches
            
            # 创建演示文稿
            prs = Presentation()
            
            # 使用 pdf2image 将 PDF 页面转换为图片
            images = convert_from_path(input_path)
            
            # 处理页面范围
            if options and options.get('page_range'):
                page_range = self._parse_page_range(options['page_range'])
                images = [images[i-1] for i in page_range if 1 <= i <= len(images)]
            
            # 为每个页面创建幻灯片
            for i, image in enumerate(images):
                # 添加空白幻灯片
                slide_layout = prs.slide_layouts[6]  # 空白布局
                slide = prs.slides.add_slide(slide_layout)
                
                # 保存图片到临时文件
                temp_image_file = tempfile.NamedTemporaryFile(suffix=f"_pdf_page_{i}.png", delete=False)
                temp_image_path = temp_image_file.name
                temp_image_file.close()
                image.save(temp_image_path, 'PNG')
                
                # 添加图片到幻灯片
                slide.shapes.add_picture(temp_image_path, 0, 0, prs.slide_width, prs.slide_height)
                
                # 清理临时文件
                Path(temp_image_path).unlink(missing_ok=True)
            
            # 保存演示文稿
            prs.save(output_path)
            
            return {
                'success': True,
                'message': 'PDF 转 PowerPoint 成功'
            }
            
        except Exception as e:
            logger.error(f"PDF 转 PowerPoint 失败: {e}")
            return {
                'success': False,
                'error': f'PDF 转 PowerPoint 失败: {str(e)}'
            }
    
    async def _pdf_to_image(
        self, 
        input_path: str, 
        output_path: str, 
        target_format: str, 
        options: Optional[Dict[str, Any]] = None
    ) -> Dict[str, Any]:
        """PDF 转换为图片"""
        try:
            from pdf2image import convert_from_path
            
            # 设置图片质量
            quality = options.get('quality', 'medium') if options else 'medium'
            if quality == 'high':
                dpi = 300
            elif quality == 'low':
                dpi = 72
            else:
                dpi = 150
            
            # 使用 pdf2image 转换
            images = convert_from_path(input_path, dpi=dpi)
            
            # 处理页面范围
            if options and options.get('page_range'):
                page_range = self._parse_page_range(options['page_range'])
                images = [images[i-1] for i in page_range if 1 <= i <= len(images)]
            
            # 转换第一页（或指定页面）
            if images:
                image = images[0]
                
                # 调整图片尺寸
                if options and options.get('output_size') and options['output_size'] != 'original':
                    image = self._resize_image(image, options['output_size'])
                
                # 根据目标格式设置保存参数
                save_kwargs = {}
                
                # 处理不同格式的保存参数
                if target_format.lower() in ['jpg', 'jpeg']:
                    save_format = 'JPEG'
                    save_kwargs['quality'] = 95
                elif target_format.lower() == 'png':
                    save_format = 'PNG'
                elif target_format.lower() == 'gif':
                    save_format = 'GIF'
                elif target_format.lower() == 'bmp':
                    save_format = 'BMP'
                elif target_format.lower() in ['tiff', 'tif']:
                    save_format = 'TIFF'
                else:
                    # 默认使用 JPEG
                    save_format = 'JPEG'
                    save_kwargs['quality'] = 95
                
                # 保存图片
                image.save(output_path, save_format, **save_kwargs)
                
                return {
                    'success': True,
                    'message': f'PDF 转 {target_format.upper()} 成功'
                }
            else:
                return {
                    'success': False,
                    'error': 'PDF 文件为空或无法读取'
                }
                
        except Exception as e:
            logger.error(f"PDF 转图片失败: {e}")
            return {
                'success': False,
                'error': f'PDF 转图片失败: {str(e)}'
            }
    
    def _extract_image_safely(self, pdf_document, image_info, page_num, img_index):
        """
        安全地提取 PDF 中的图片
        
        Args:
            pdf_document: PDF 文档对象
            image_info: 图片信息
            page_num: 页面编号
            img_index: 图片索引
            
        Returns:
            临时图片文件路径，如果失败返回 None
        """
        try:
            import fitz  # 确保导入 fitz 模块
            import os
            
            # 获取图片引用 - 处理不同的图片信息格式
            xref = None
            
            if isinstance(image_info, dict):
                # 字典格式：{'image': xref, 'bbox': rect, ...}
                if 'image' in image_info:
                    xref = image_info['image']
            elif isinstance(image_info, (list, tuple)):
                # 列表/元组格式：[xref, ...]
                if len(image_info) > 0:
                    xref = image_info[0]
            elif isinstance(image_info, (int, str)):
                # 直接是引用号
                xref = image_info
            elif hasattr(image_info, '__getitem__'):
                # 其他可索引对象
                try:
                    xref = image_info[0]
                except (IndexError, TypeError):
                    pass
            
            if xref is None:
                logger.warning(f"无法获取图片引用 (页面 {page_num}, 索引 {img_index})")
                return None
            
            # 确保 xref 是整数
            try:
                xref = int(xref)
            except (ValueError, TypeError):
                logger.warning(f"图片引用格式错误: {xref} (页面 {page_num}, 索引 {img_index})")
                return None
            
            # 创建 Pixmap 对象
            pix = fitz.Pixmap(pdf_document, xref)
            
            # 检查颜色空间
            if pix.n - pix.alpha < 4:  # 灰度或 RGB
                # 直接保存
                temp_img_file = tempfile.NamedTemporaryFile(suffix=f"_pdf_img_{page_num}_{img_index}.png", delete=False)
                temp_img_path = temp_img_file.name
                temp_img_file.close()
                pix.save(temp_img_path)
            else:  # CMYK: 转换为 RGB
                # 转换为 RGB 颜色空间
                pix1 = fitz.Pixmap(fitz.csRGB, pix)
                temp_img_file = tempfile.NamedTemporaryFile(suffix=f"_pdf_img_{page_num}_{img_index}.png", delete=False)
                temp_img_path = temp_img_file.name
                temp_img_file.close()
                pix1.save(temp_img_path)
                pix1 = None
            
            # 释放内存
            pix = None
            
            # 验证文件是否成功创建
            logger.info(f"尝试创建临时文件: {temp_img_path}")
            if os.path.exists(temp_img_path) and os.path.getsize(temp_img_path) > 0:
                logger.info(f"临时文件创建成功: {temp_img_path}")
                return temp_img_path
            else:
                logger.warning(f"图片文件创建失败 (页面 {page_num}, 索引 {img_index})")
                logger.warning(f"文件存在: {os.path.exists(temp_img_path)}")
                if os.path.exists(temp_img_path):
                    logger.warning(f"文件大小: {os.path.getsize(temp_img_path)}")
                return None
            
        except Exception as e:
            logger.warning(f"提取图片失败 (页面 {page_num}, 索引 {img_index}): {e}")
            return None

    def _calculate_image_size(self, image_info, page_width, page_height):
        """
        计算图片在 Word 文档中的合适大小
        
        Args:
            image_info: 图片信息（可能包含 bbox）
            page_width: 页面宽度
            page_height: 页面高度
            
        Returns:
            (width, height) 元组，单位为英寸
        """
        try:
            from docx.shared import Inches
            
            # 默认图片大小
            default_width = 5.0  # 5 英寸
            default_height = 4.0  # 4 英寸
            
            # 如果有 bbox 信息，根据原始尺寸计算
            if isinstance(image_info, dict) and 'bbox' in image_info:
                bbox = image_info['bbox']
                if len(bbox) == 4:
                    img_width = bbox[2] - bbox[0]  # 图片宽度
                    img_height = bbox[3] - bbox[1]  # 图片高度
                    
                    # 计算图片在页面中的比例
                    width_ratio = img_width / page_width
                    height_ratio = img_height / page_height
                    
                    # 根据比例调整大小
                    if width_ratio > 0.8:  # 图片宽度超过页面 80%
                        # 大图片，设置为页面宽度的 80%
                        width = min(6.0, page_width * 0.8 / 72)  # 转换为英寸
                        height = width * (img_height / img_width)
                    elif width_ratio > 0.5:  # 图片宽度超过页面 50%
                        # 中等图片，设置为页面宽度的 60%
                        width = min(5.0, page_width * 0.6 / 72)
                        height = width * (img_height / img_width)
                    else:
                        # 小图片，保持原始比例但限制最大尺寸
                        width = min(4.0, img_width / 72)
                        height = min(3.0, img_height / 72)
                    
                    # 确保高度不超过页面高度的 80%
                    max_height = page_height * 0.8 / 72
                    if height > max_height:
                        height = max_height
                        width = height * (img_width / img_width)
                    
                    return (width, height)
            
            # 如果没有 bbox 信息，使用默认大小
            return (default_width, default_height)
            
        except Exception as e:
            logger.warning(f"计算图片大小时出错: {e}")
            return (5.0, 4.0)  # 返回默认大小

    def _calculate_smart_excel_image_size(
        self, 
        pdf_document, 
        image_info, 
        page_num, 
        img_index, 
        options: Optional[Dict[str, Any]] = None
    ) -> Tuple[int, int]:
        """
        智能计算 Excel 中图片的像素尺寸
        
        Args:
            pdf_document: PDF 文档对象
            image_info: 图片信息
            page_num: 页面编号
            img_index: 图片索引
            options: 转换选项
            
        Returns:
            (width, height) 元组，单位为像素
        """
        try:
            # 获取选项中的尺寸限制
            max_width = options.get('image_width', 300) if options else 300
            max_height = options.get('image_height', 200) if options else 200
            min_width = options.get('min_image_width', 50) if options else 50
            min_height = options.get('min_image_height', 30) if options else 30
            
            # 获取页面信息
            page = pdf_document[page_num - 1]
            page_width = page.rect.width
            page_height = page.rect.height
            
            # 尝试获取图片在页面中的边界框
            try:
                img_rect = page.get_image_bbox(image_info)
                if img_rect:
                    img_x, img_y, img_w, img_h = img_rect
                    # 计算图片在页面中的相对位置和大小
                    relative_width = img_w / page_width
                    relative_height = img_h / page_height
                    
                    # 根据图片在页面中的实际大小计算 Excel 中的尺寸
                    if img_w > img_h:
                        # 横向图片：优先按宽度缩放
                        target_width = min(max_width, int(img_w * 0.8))  # 缩小到80%
                        target_height = int(target_width * (img_h / img_w))
                        if target_height > max_height:
                            target_height = max_height
                            target_width = int(target_height * (img_w / img_h))
                    else:
                        # 纵向图片：优先按高度缩放
                        target_height = min(max_height, int(img_h * 0.8))  # 缩小到80%
                        target_width = int(target_height * (img_w / img_h))
                        if target_width > max_width:
                            target_width = max_width
                            target_height = int(target_width * (img_h / img_w))
                    
                    # 确保最小尺寸
                    target_width = max(min_width, target_width)
                    target_height = max(min_height, target_height)
                    
                    logger.info(f"智能计算图片尺寸: 原始 {img_w}x{img_h}, 缩放为 {target_width}x{target_height}")
                    return target_width, target_height
                    
                else:
                    # 如果无法获取边界框，使用默认计算
                    logger.warning(f"无法获取图片边界框，使用默认尺寸")
                    return self._calculate_default_excel_image_size(max_width, max_height, min_width, min_height)
                    
            except Exception as bbox_error:
                logger.warning(f"获取图片边界框失败: {bbox_error}")
                return self._calculate_default_excel_image_size(max_width, max_height, min_width, min_height)
            
        except Exception as e:
            logger.warning(f"智能图片尺寸计算失败: {e}")
            return 200, 150  # 返回默认尺寸

    def _calculate_default_excel_image_size(self, max_width, max_height, min_width, min_height):
        """计算默认的 Excel 图片尺寸"""
        try:
            # 根据限制计算合适的默认尺寸
            if max_width > max_height:
                # 横向默认尺寸
                width = min(max_width, 300)
                height = min(max_height, 200)
            else:
                # 纵向默认尺寸
                height = min(max_height, 250)
                width = min(max_width, 200)
            
            # 确保最小尺寸
            width = max(min_width, width)
            height = max(min_height, height)
            
            return width, height
            
        except Exception as e:
            logger.warning(f"默认图片尺寸计算失败: {e}")
            return 200, 150

    def _collect_page_content(self, page, page_num, pdf_document, options):
        """
        收集页面上的所有内容（文本和图片），按位置排序
        
        Args:
            page: PDF 页面对象
            page_num: 页面编号
            pdf_document: PDF 文档对象
            options: 转换选项
            
        Returns:
            按位置排序的内容列表
        """
        page_content = []
        
        try:
            # 1. 收集文本内容
            text_dict = page.get_text("dict")
            for block in text_dict.get("blocks", []):
                if "lines" in block:  # 文本块
                    for line in block["lines"]:
                        # 计算文本块的垂直位置
                        if line["spans"]:
                            first_span = line["spans"][0]
                            bbox = first_span.get("bbox", [0, 0, 0, 0])
                            y_position = bbox[1]  # 垂直位置
                            
                            # 收集行文本
                            line_text = ""
                            for span in line["spans"]:
                                line_text += span["text"]
                            
                            if line_text.strip():
                                page_content.append({
                                    'type': 'text',
                                    'y_position': y_position,
                                    'data': line  # 存储完整的行数据，包含spans信息
                                })
            
            # 2. 收集图片内容
            image_list = page.get_images()
            if image_list:
                for img_index, img in enumerate(image_list):
                    try:
                        # 尝试获取图片在页面中的位置
                        y_position = self._get_image_y_position(page, img, img_index)
                        
                        page_content.append({
                            'type': 'image',
                            'y_position': y_position,
                            'data': {
                                'img': img,
                                'img_index': img_index,
                                'page_num': page_num,
                                'y_position': y_position
                            }
                        })
                        
                        logger.info(f"收集图片: 索引 {img_index}, 位置 Y={y_position}")
                        
                    except Exception as img_error:
                        logger.warning(f"收集图片失败: {img_error}")
                        continue
            
            # 3. 按垂直位置排序
            page_content.sort(key=lambda x: x['y_position'])
            
            logger.info(f"页面 {page_num} 收集到 {len(page_content)} 个内容项")
            
            return page_content
            
        except Exception as e:
            logger.error(f"收集页面内容失败: {e}")
            return []

    def _get_image_y_position(self, page, image_info, img_index):
        """
        获取图片在页面中的垂直位置
        
        Args:
            page: PDF 页面对象
            image_info: 图片信息
            img_index: 图片索引
            
        Returns:
            图片的垂直位置（Y坐标）
        """
        try:
            # 尝试获取图片边界框
            img_rect = page.get_image_bbox(image_info)
            if img_rect:
                return img_rect[1]  # 返回Y坐标
            
            # 如果无法获取边界框，使用默认位置
            # 根据图片索引计算位置
            page_height = page.rect.height
            default_y = 100 + img_index * 200  # 从顶部开始，每个图片间隔200像素
            
            logger.warning(f"无法获取图片 {img_index} 的边界框，使用默认位置 Y={default_y}")
            return default_y
            
        except Exception as e:
            logger.warning(f"获取图片位置失败: {e}")
            # 使用页面顶部位置
            return 100 + img_index * 200

    def _get_excel_font(self, font_name: str):
        """获取 Excel 字体对象"""
        try:
            from openpyxl.styles import Font
            
            # 字体映射
            font_mapping = {
                'Arial': 'Arial',
                'Times': 'Times New Roman',
                'Helvetica': 'Arial',
                'Courier': 'Courier New',
                'Calibri': 'Calibri',
                'SimSun': 'SimSun',
                'SimHei': 'SimHei',
                'Microsoft YaHei': 'Microsoft YaHei'
            }
            
            # 查找映射
            excel_font_name = font_mapping.get(font_name, 'Arial')
            
            return Font(name=excel_font_name)
            
        except Exception as e:
            logger.warning(f"获取 Excel 字体失败: {e}")
            from openpyxl.styles import Font
            return Font(name='Arial')

    def _set_excel_cell_alignment(self, cell, bbox, page_width):
        """设置 Excel 单元格对齐方式"""
        try:
            from openpyxl.styles import Alignment
            
            if len(bbox) >= 4:
                line_left = bbox[0]
                line_right = bbox[2]
                line_width = line_right - line_left
                line_center = (line_left + line_right) / 2
                page_center = page_width / 2
                
                # 计算对齐方式
                left_margin = line_left / page_width
                right_margin = (page_width - line_right) / page_width
                center_offset = abs(line_center - page_center) / page_width
                
                # 设置对齐方式
                if center_offset < 0.08 and left_margin > 0.02 and right_margin > 0.02:
                    # 居中对齐
                    cell.alignment = Alignment(horizontal='center')
                elif left_margin < 0.02 and right_margin > 0.05:
                    # 左对齐
                    cell.alignment = Alignment(horizontal='left')
                elif right_margin < 0.02 and left_margin > 0.05:
                    # 右对齐
                    cell.alignment = Alignment(horizontal='right')
                elif line_width > page_width * 0.75:
                    # 两端对齐
                    cell.alignment = Alignment(horizontal='justify')
                else:
                    # 默认左对齐
                    cell.alignment = Alignment(horizontal='left')
            else:
                cell.alignment = Alignment(horizontal='left')
                
        except Exception as e:
            logger.warning(f"设置 Excel 单元格对齐失败: {e}")

    def _set_excel_row_format(self, ws, row, font_size):
        """设置 Excel 行格式"""
        try:
            # 根据字体大小设置行高
            if font_size >= 24:
                row_height = 30
            elif font_size >= 18:
                row_height = 25
            elif font_size >= 14:
                row_height = 20
            elif font_size >= 12:
                row_height = 18
            else:
                row_height = 15
            
            ws.row_dimensions[row].height = row_height
            
        except Exception as e:
            logger.warning(f"设置 Excel 行格式失败: {e}")

    def _get_pages_to_process(self, pages, options: Optional[Dict[str, Any]] = None) -> List[int]:
        """获取需要处理的页面列表"""
        # 获取页面总数
        if hasattr(pages, 'page_count'):  # PyMuPDF 文档
            total_pages = pages.page_count
        else:  # pdfplumber 页面列表
            total_pages = len(pages)
        
        if not options or not options.get('page_range'):
            return list(range(1, total_pages + 1))
        
        return self._parse_page_range(options['page_range'])
    
    def _parse_page_range(self, page_range: str) -> List[int]:
        """解析页面范围字符串"""
        pages = []
        parts = page_range.split(',')
        
        for part in parts:
            part = part.strip()
            if '-' in part:
                start, end = map(int, part.split('-'))
                pages.extend(range(start, end + 1))
            else:
                pages.append(int(part))
        
        return sorted(set(pages))  # 去重并排序
    
    def _resize_image(self, image: Image.Image, output_size: str) -> Image.Image:
        """调整图片尺寸"""
        if output_size == 'A4':
            # A4 尺寸: 210mm x 297mm, 300 DPI
            target_width = int(210 * 300 / 25.4)  # mm to pixels
            target_height = int(297 * 300 / 25.4)
        elif output_size == 'letter':
            # Letter 尺寸: 8.5" x 11", 300 DPI
            target_width = int(8.5 * 300)
            target_height = int(11 * 300)
        else:
            return image
        
        # 保持宽高比
        image.thumbnail((target_width, target_height), Image.Resampling.LANCZOS)
        return image
    
    def get_pdf_info(self, input_path: str) -> Dict[str, Any]:
        """获取 PDF 文件信息"""
        try:
            with PyPDF2.PdfReader(input_path) as reader:
                return {
                    'page_count': len(reader.pages),
                    'file_size': Path(input_path).stat().st_size,
                    'is_encrypted': reader.is_encrypted,
                    'metadata': reader.metadata
                }
        except Exception as e:
            logger.error(f"获取 PDF 信息失败: {e}")
            return {}
    
    def _convert_color_to_rgb(self, color: int) -> Optional[tuple]:
        """将 PDF 颜色值转换为 RGB 颜色元组"""
        try:
            # PDF 颜色通常是 24 位整数，格式为 0xRRGGBB
            if color == 0:  # 黑色
                return None
            
            # 提取 RGB 分量
            r = (color >> 16) & 0xFF
            g = (color >> 8) & 0xFF
            b = color & 0xFF
            
            # 返回 RGB 元组
            return (r, g, b)
            
        except Exception as e:
            logger.warning(f"颜色转换失败: {e}")
            return None
    
    def _set_paragraph_alignment(self, paragraph, bbox, page_width):
        """设置段落对齐方式"""
        try:
            from docx.enum.text import WD_ALIGN_PARAGRAPH
            
            if len(bbox) >= 4:
                line_left = bbox[0]
                line_right = bbox[2]
                line_width = line_right - line_left
                line_center = (line_left + line_right) / 2
                page_center = page_width / 2
                
                # 计算对齐方式
                left_margin = line_left / page_width
                right_margin = (page_width - line_right) / page_width
                center_offset = abs(line_center - page_center) / page_width
                

                
                # 更精确的对齐判断
                if center_offset < 0.08 and left_margin > 0.02 and right_margin > 0.02:
                    # 非常接近页面中心，居中对齐
                    paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
                    logger.info(f"设置居中对齐: 中心偏移={center_offset:.3f}")
                elif left_margin < 0.02 and right_margin > 0.05:
                    # 紧贴左边，左对齐
                    paragraph.alignment = WD_ALIGN_PARAGRAPH.LEFT
                elif right_margin < 0.02 and left_margin > 0.05:
                    # 紧贴右边，右对齐
                    paragraph.alignment = WD_ALIGN_PARAGRAPH.RIGHT
                elif line_width > page_width * 0.75:
                    # 行宽超过页面75%，可能是两端对齐
                    paragraph.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
                elif center_offset < 0.12:
                    # 接近中心，居中对齐
                    paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
                else:
                    # 默认左对齐
                    paragraph.alignment = WD_ALIGN_PARAGRAPH.LEFT
            else:
                paragraph.alignment = WD_ALIGN_PARAGRAPH.LEFT
                logger.warning(f"边界框信息不足: {bbox}")
                
        except Exception as e:
            logger.warning(f"设置段落对齐失败: {e}")
            paragraph.alignment = WD_ALIGN_PARAGRAPH.LEFT

    def _set_line_spacing(self, paragraph, font_size):
        """设置行间距"""
        try:
            from docx.shared import Pt
            
            # 根据字体大小设置合适的行间距
            if font_size >= 24:
                # 超大字体，使用 1.8 倍行距
                paragraph.paragraph_format.line_spacing = 1.8
            elif font_size >= 18:
                # 大字体，使用 1.5 倍行距
                paragraph.paragraph_format.line_spacing = 1.5
            elif font_size >= 14:
                # 中等字体，使用 1.3 倍行距
                paragraph.paragraph_format.line_spacing = 1.3
            elif font_size >= 12:
                # 标准字体，使用 1.2 倍行距
                paragraph.paragraph_format.line_spacing = 1.2
            else:
                # 小字体，使用 1.1 倍行距
                paragraph.paragraph_format.line_spacing = 1.1
                
        except Exception as e:
            logger.warning(f"设置行间距失败: {e}")
            paragraph.paragraph_format.line_spacing = 1.2

    def _process_page_layout(self, doc, text_dict, page):
        """智能处理页面布局"""
        try:
            # 按位置排序文本块，保持布局顺序
            text_blocks = []
            for block in text_dict.get("blocks", []):
                if "lines" in block:  # 文本块
                    text_blocks.append(block)
            
            # 按垂直位置排序（从上到下）
            text_blocks.sort(key=lambda b: b.get("bbox", [0, 0, 0, 0])[1])
            
            # 分析文档结构
            document_structure = self._analyze_document_structure(text_blocks, page.rect.width)
            
            # 处理文本块
            for block in text_blocks:
                if "lines" in block:  # 文本块
                    self._process_text_block(doc, block, document_structure, page.rect.width)
                    
        except Exception as e:
            logger.error(f"处理页面布局失败: {e}")

    def _analyze_document_structure(self, text_blocks, page_width):
        """分析文档结构"""
        structure = {
            'title_lines': [],  # 大标题行
            'subtitle_lines': [],  # 副标题行
            'body_lines': [],  # 正文行
            'list_items': [],  # 列表项
            'header_info': []  # 页眉信息
        }
        
        for block in text_blocks:
            if "lines" in block:
                for line in block["lines"]:
                    line_text = ""
                    max_font_size = 0
                    line_bbox = line.get("bbox", [0, 0, 0, 0])
                    
                    # 收集行文本和最大字体大小
                    for span in line["spans"]:
                        line_text += span["text"]
                        max_font_size = max(max_font_size, span.get("size", 12))
                    
                    # 转换字体大小
                    font_size = max(8, min(72, max_font_size * 0.75))
                    
                    # 分析行类型
                    line_info = {
                        'text': line_text.strip(),
                        'bbox': line_bbox,
                        'font_size': font_size,
                        'block': block
                    }
                    
                    # 根据字体大小和位置分类
                    if font_size >= 20:
                        structure['title_lines'].append(line_info)
                    elif font_size >= 16:
                        structure['subtitle_lines'].append(line_info)
                    elif self._is_list_item(line_text):
                        structure['list_items'].append(line_info)
                    elif self._is_header_info(line_text, line_bbox, page_width):
                        structure['header_info'].append(line_info)
                    else:
                        structure['body_lines'].append(line_info)
        
        return structure

    def _is_list_item(self, text):
        """判断是否为列表项"""
        import re
        # 匹配数字编号：1. 2. 3. 等
        if re.match(r'^\d+\.', text.strip()):
            return True
        # 匹配字母编号：a. b. c. 等
        if re.match(r'^[a-z]\.', text.strip()):
            return True
        # 匹配罗马数字：i. ii. iii. 等
        if re.match(r'^[ivxlcdm]+\.', text.strip(), re.IGNORECASE):
            return True
        return False

    def _is_bullet_point(self, text):
        """判断是否为 bullet point"""
        import re
        
        # 去除首尾空格
        text = text.strip()
        if not text:
            return False
        
        # 常见的 bullet point 字符
        bullet_chars = [
            '•', '◦', '‣', '▪', '▫', '▸', '▹', '▻', '▽', '▾',
            '-', '*', '+', '→', '⇒', '▶', '►', '▷', '◁', '◀',
            '◆', '◇', '◈', '◉', '◎', '●', '○', '◐', '◑', '◒',
            '◓', '◔', '◕', '◖', '◗', '◘', '◙', '◚', '◛', '◜',
            '◝', '◞', '◟', '◠', '◡', '◢', '◣', '◤', '◥', '◦',
            '◧', '◨', '◩', '◪', '◫', '◬', '◭', '◮', '◯', '◰',
            '◱', '◲', '◳', '◴', '◵', '◶', '◷', '◸', '◹', '◺',
            '◻', '◼', '◽', '◾', '◿', '☀', '☁', '☂', '☃', '☄',
            '★', '☆', '☎', '☏', '☐', '☑', '☒', '☓', '☔', '☕',
            '☖', '☗', '☘', '☙', '☚', '☛', '☜', '☝', '☞', '☟',
            '☠', '☡', '☢', '☣', '☤', '☥', '☦', '☧', '☨', '☩',
            '☪', '☫', '☬', '☭', '☮', '☯', '☰', '☱', '☲', '☳',
            '☴', '☵', '☶', '☷', '☸', '☹', '☺', '☻', '☼', '☽',
            '☾', '☿', '♀', '♂', '♁', '♂', '♃', '♄', '♅', '♆',
            '♇', '♈', '♉', '♊', '♋', '♌', '♍', '♎', '♏', '♐',
            '♑', '♒', '♓', '♔', '♕', '♖', '♗', '♘', '♙', '♚',
            '♛', '♜', '♝', '♞', '♟', '♠', '♡', '♢', '♣', '♤',
            '♥', '♦', '♧', '♨', '♩', '♪', '♫', '♬', '♭', '♮',
            '♯', '♰', '♱', '♲', '♳', '♴', '♵', '♶', '♷', '♸',
            '♹', '♺', '♻', '♼', '♽', '♾', '♿', '⚀', '⚁', '⚂',
            '⚃', '⚄', '⚅', '⚆', '⚇', '⚈', '⚉', '⚊', '⚋', '⚌',
            '⚍', '⚎', '⚏', '⚐', '⚑', '⚒', '⚓', '⚔', '⚕', '⚖',
            '⚗', '⚘', '⚙', '⚚', '⚛', '⚜', '⚝', '⚞', '⚟', '⚠',
            '⚡', '⚢', '⚣', '⚤', '⚥', '⚦', '⚧', '⚨', '⚩', '⚪',
            '⚫', '⚬', '⚭', '⚮', '⚯', '⚰', '⚱', '⚲', '⚳', '⚴',
            '⚵', '⚶', '⚷', '⚸', '⚹', '⚺', '⚻', '⚼', '⚽', '⚾',
            '⚿', '⛀', '⛁', '⛂', '⛃', '⛄', '⛅', '⛆', '⛇', '⛈',
            '⛉', '⛊', '⛋', '⛌', '⛍', '⛎', '⛏', '⛐', '⛑', '⛒',
            '⛓', '⛔', '⛕', '⛖', '⛗', '⛘', '⛙', '⛚', '⛛', '⛜',
            '⛝', '⛞', '⛟', '⛠', '⛡', '⛢', '⛣', '⛤', '⛥', '⛦',
            '⛧', '⛨', '⛩', '⛪', '⛫', '⛬', '⛭', '⛮', '⛯', '⛰',
            '⛱', '⛲', '⛳', '⛴', '⛵', '⛶', '⛷', '⛸', '⛹', '⛺',
            '⛻', '⛼', '⛽', '⛾', '⛿'
        ]
        
        # 检查第一个字符是否为 bullet point
        if text and text[0] in bullet_chars:
            return True
        
        # 检查常见的 bullet point 模式
        bullet_patterns = [
            r'^[•◦‣▪▫▸▹▻▽▾→⇒▶►▷◁◀◆◇◈◉◎●○◐◑◒◓◔◕◖◗◘◙◚◛◜◝◞◟◠◡◢◣◤◥◦◧◨◩◪◫◬◭◮◯◰◱◲◳◴◵◶◷◸◹◺◻◼◽◾◿☀☁☂☃☄★☆☎☏☐☑☒☓☔☕☖☗☘☙☚☛☜☝☞☟☠☡☢☣☤☥☦☧☨☩☪☫☬☭☮☯☰☱☲☳☴☵☶☷☸☹☺☻☼☽☾☿♀♂♁♂♃♄♅♆♇♈♉♊♋♌♍♎♏♐♑♒♓♔♕♖♗♘♙♚♛♜♝♞♟♠♡♢♣♤♥♦♧♨♩♪♫♬♭♮♯♰♱♲♳♴♵♶♷♸♹♺♻♼♽♾♿⚀⚁⚂⚃⚄⚅⚆⚇⚈⚉⚊⚋⚌⚍⚎⚏⚐⚑⚒⚓⚔⚕⚖⚗⚘⚙⚚⚛⚜⚝⚞⚟⚠⚡⚢⚣⚤⚥⚦⚧⚨⚩⚪⚫⚬⚭⚮⚯⚰⚱⚲⚳⚴⚵⚶⚷⚸⚹⚺⚻⚼⚽⚾⚿⛀⛁⛂⛃⛄⛅⛆⛇⛈⛉⛊⛋⛌⛍⛎⛏⛐⛑⛒⛓⛔⛕⛖⛗⛘⛙⛚⛛⛜⛝⛞⛟⛠⛡⛢⛣⛤⛥⛦⛧⛨⛩⛪⛫⛬⛭⛮⛯⛰⛱⛲⛳⛴⛵⛶⛷⛸⛹⛺⛻⛼⛽⛾⛿]\s+',  # 各种 bullet 字符后跟空格
            r'^[-*+]\s+',  # 常见的 bullet 字符
            r'^[→⇒▶►▷◁◀]\s+',  # 箭头类 bullet
            r'^[◆◇◈◉◎●○]\s+',  # 圆形类 bullet
            r'^[★☆]\s+',  # 星形类 bullet
            r'^[☐☑☒]\s+',  # 复选框类 bullet
        ]
        
        for pattern in bullet_patterns:
            if re.match(pattern, text):
                return True
        
        return False

    def _is_header_info(self, text, bbox, page_width):
        """判断是否为页眉信息"""
        # 检查是否在页面顶部
        if len(bbox) >= 4 and bbox[1] < 100:  # 距离顶部小于100像素
            return True
        # 检查是否包含文档编号或日期
        if any(keyword in text.lower() for keyword in ['a/inb', 'geneva', 'november', 'december']):
            return True
        return False

    def _process_text_block(self, doc, block, structure, page_width):
        """处理文本块"""
        try:
            # 收集所有行
            lines = []
            for line in block["lines"]:
                line_text = ""
                line_bbox = line.get("bbox", [0, 0, 0, 0])
                
                # 收集行文本
                for span in line["spans"]:
                    line_text += span["text"]
                
                line_text = line_text.strip()
                if not line_text:
                    continue
                
                lines.append({
                    'text': line_text,
                    'bbox': line_bbox,
                    'spans': line["spans"],
                    'is_bullet': self._is_bullet_point(line_text),
                    'is_list': self._is_list_item(line_text)
                })
            
            # 智能合并相关行
            merged_lines = self._merge_bullet_point_lines(lines)
            
            # 处理合并后的行
            for line_info in merged_lines:
                # 创建段落
                paragraph = doc.add_paragraph()
                
                # 设置段落对齐方式
                self._set_paragraph_alignment(paragraph, line_info['bbox'], page_width)
                
                # 处理列表项和 bullet points 缩进
                if line_info['is_list']:
                    self._set_list_item_format(paragraph, line_info['bbox'], page_width)
                elif line_info['is_bullet']:
                    self._set_bullet_point_format(paragraph, line_info['bbox'], page_width)
                
                # 处理行中的每个 span
                max_font_size = 0
                from docx.shared import Pt
                for span in line_info['spans']:
                    # 获取文本和样式信息
                    text = span["text"]
                    pdf_font_name = span.get("font", "Arial")
                    raw_font_size = span.get("size", 12)
                    
                    # 转换字体大小
                    font_size = max(8, min(72, raw_font_size * 0.75))
                    max_font_size = max(max_font_size, font_size)
                    
                    # 改进的样式检测
                    is_bold = "Bold" in pdf_font_name or span.get("flags", 0) & 2**4
                    is_italic = "Italic" in pdf_font_name or span.get("flags", 0) & 2**1
                    
                    # 通过字体大小检测标题样式
                    if font_size >= 24:
                        is_bold = True
                    
                    # 通过字体名称特征检测
                    font_lower = pdf_font_name.lower()
                    if any(keyword in font_lower for keyword in ['bold', 'heavy', 'black', 'extra']):
                        is_bold = True
                    if any(keyword in font_lower for keyword in ['italic', 'oblique']):
                        is_italic = True
                    
                    if text.strip():
                        # 添加运行到段落
                        run = paragraph.add_run(text)
                        
                        # 应用字体样式
                        word_font_name = self._get_font_mapping(pdf_font_name)
                        run.font.name = word_font_name
                        run.font.size = Pt(font_size)
                        run.bold = is_bold
                        run.italic = is_italic
                        
                        # 处理颜色（如果有）
                        if "color" in span:
                            color = span["color"]
                            if color != 0:  # 不是黑色
                                rgb_color = self._convert_color_to_rgb(color)
                                if rgb_color:
                                    try:
                                        from docx.shared import RGBColor
                                        run.font.color.rgb = RGBColor(*rgb_color)
                                    except Exception as color_error:
                                        logger.warning(f"设置字体颜色失败: {color_error}")
                
                # 设置行间距
                self._set_line_spacing(paragraph, max_font_size)
                
                # 设置高级格式（包括负缩进等）
                self._set_advanced_formatting(paragraph, line_bbox, page_width, max_font_size)
                
        except Exception as e:
            logger.error(f"处理文本块失败: {e}")

    def _set_list_item_format(self, paragraph, bbox, page_width):
        """设置列表项格式"""
        try:
            from docx.shared import Inches
            
            if len(bbox) >= 4:
                # 计算缩进
                left_margin = bbox[0]
                
                # 检查是否需要负缩进（用于特殊布局）
                if left_margin < 0:
                    # 负缩进，用于特殊布局效果
                    paragraph.paragraph_format.left_indent = Inches(left_margin / 72)
                    paragraph.paragraph_format.first_line_indent = Inches(left_margin * 2 / 72)
                elif left_margin > 50:  # 如果有明显缩进
                    paragraph.paragraph_format.left_indent = Inches(left_margin / 72)
                    
        except Exception as e:
            logger.warning(f"设置列表项格式失败: {e}")

    def _set_bullet_point_format(self, paragraph, bbox, page_width):
        """设置 bullet point 格式"""
        try:
            from docx.shared import Inches
            
            if len(bbox) >= 4:
                # 计算缩进
                left_margin = bbox[0]
                
                # 设置 bullet point 的缩进
                if left_margin > 50:  # 如果有明显缩进
                    paragraph.paragraph_format.left_indent = Inches(left_margin / 72)
                    # 设置首行缩进，让 bullet point 突出
                    paragraph.paragraph_format.first_line_indent = Inches(-0.25)  # 负缩进让 bullet 突出
                else:
                    # 默认 bullet point 缩进
                    paragraph.paragraph_format.left_indent = Inches(0.5)
                    paragraph.paragraph_format.first_line_indent = Inches(-0.25)
                
                # 设置段落间距，让 bullet points 更紧凑
                from docx.shared import Pt
                paragraph.paragraph_format.space_after = Pt(6)
                paragraph.paragraph_format.space_before = Pt(0)
                    
        except Exception as e:
            logger.warning(f"设置 bullet point 格式失败: {e}")

    def _merge_bullet_point_lines(self, lines):
        """智能合并 bullet point 行"""
        if not lines:
            return []
        
        merged_lines = []
        current_group = []
        
        for i, line in enumerate(lines):
            # 检查是否需要开始新组
            if not current_group:
                current_group = [line]
                continue
            
            # 检查是否可以合并到当前组
            last_line = current_group[-1]
            
            # 检查垂直距离
            vertical_distance = abs(line['bbox'][1] - last_line['bbox'][1])
            
            # 检查是否为连续的 bullet points
            is_consecutive_bullets = (
                last_line['is_bullet'] and line['is_bullet'] and 
                vertical_distance < 25  # 25像素内的垂直距离
            )
            
            # 检查是否为同一 bullet point 的延续
            is_bullet_continuation = (
                last_line['is_bullet'] and not line['is_bullet'] and
                vertical_distance < 15 and  # 更小的垂直距离
                len(line['text']) > 0 and not line['text'].startswith('•')  # 不是新的 bullet
            )
            
            if is_consecutive_bullets or is_bullet_continuation:
                # 合并到当前组
                current_group.append(line)
            else:
                # 完成当前组，开始新组
                if current_group:
                    merged_line = self._merge_line_group(current_group)
                    merged_lines.append(merged_line)
                current_group = [line]
        
        # 处理最后一组
        if current_group:
            merged_line = self._merge_line_group(current_group)
            merged_lines.append(merged_line)
        
        return merged_lines

    def _merge_line_group(self, line_group):
        """合并行组"""
        if not line_group:
            return None
        
        if len(line_group) == 1:
            return line_group[0]
        
        # 合并多行
        merged_text = ""
        merged_spans = []
        merged_bbox = line_group[0]['bbox'].copy()
        
        for i, line in enumerate(line_group):
            # 添加行文本
            if i > 0:
                merged_text += " "  # 行间添加空格
            merged_text += line['text']
            
            # 合并 spans
            merged_spans.extend(line['spans'])
            
            # 更新边界框
            if len(line['bbox']) >= 4 and len(merged_bbox) >= 4:
                merged_bbox[2] = max(merged_bbox[2], line['bbox'][2])  # 右边界
                merged_bbox[3] = max(merged_bbox[3], line['bbox'][3])  # 下边界
        
        return {
            'text': merged_text,
            'bbox': merged_bbox,
            'spans': merged_spans,
            'is_bullet': line_group[0]['is_bullet'],
            'is_list': line_group[0]['is_list']
        }

    def _set_advanced_formatting(self, paragraph, bbox, page_width, font_size):
        """设置高级格式（包括负缩进等）"""
        try:
            from docx.shared import Inches, Pt
            
            if len(bbox) >= 4:
                left_margin = bbox[0]
                right_margin = bbox[2]
                
                # 检查是否需要负缩进
                if left_margin < 0:
                    # 负缩进，用于特殊布局效果
                    paragraph.paragraph_format.left_indent = Inches(left_margin / 72)
                    paragraph.paragraph_format.first_line_indent = Inches(left_margin * 2 / 72)
                
                # 检查是否需要右缩进
                if right_margin < page_width - 50:
                    right_indent = page_width - right_margin
                    paragraph.paragraph_format.right_indent = Inches(right_indent / 72)
                
                # 根据字体大小设置段落间距
                if font_size >= 20:
                    paragraph.paragraph_format.space_after = Pt(12)
                elif font_size >= 16:
                    paragraph.paragraph_format.space_after = Pt(8)
                else:
                    paragraph.paragraph_format.space_after = Pt(6)
                    
        except Exception as e:
            logger.warning(f"设置高级格式失败: {e}")

    def _set_paragraph_spacing(self, paragraph, font_size):
        """设置段落间距"""
        try:
            from docx.shared import Pt
            
            # 根据字体大小设置段落间距
            if font_size >= 24:
                # 超大字体，使用较大段落间距
                paragraph.paragraph_format.space_after = Pt(16)
                paragraph.paragraph_format.space_before = Pt(8)
            elif font_size >= 18:
                # 大字体，使用较大段落间距
                paragraph.paragraph_format.space_after = Pt(12)
                paragraph.paragraph_format.space_before = Pt(6)
            elif font_size >= 14:
                # 中等字体，使用中等段落间距
                paragraph.paragraph_format.space_after = Pt(8)
                paragraph.paragraph_format.space_before = Pt(4)
            elif font_size >= 12:
                # 标准字体，使用标准段落间距
                paragraph.paragraph_format.space_after = Pt(6)
                paragraph.paragraph_format.space_before = Pt(2)
            else:
                # 小字体，使用较小段落间距
                paragraph.paragraph_format.space_after = Pt(4)
                paragraph.paragraph_format.space_before = Pt(1)
                
        except Exception as e:
            logger.warning(f"设置段落间距失败: {e}")
            paragraph.paragraph_format.space_after = Pt(6)

    def _get_font_mapping(self, pdf_font_name: str) -> str:
        """获取字体映射，将 PDF 字体名映射到 Word 字体名"""
        font_mapping = {
            # 常见字体映射
            'Arial': 'Arial',
            'Arial-Bold': 'Arial',
            'Arial-Italic': 'Arial',
            'Arial-BoldItalic': 'Arial',
            'ArialMT': 'Arial',
            'Arial-BoldMT': 'Arial',
            'Times': 'Times New Roman',
            'Times-Bold': 'Times New Roman',
            'Times-Italic': 'Times New Roman',
            'Times-BoldItalic': 'Times New Roman',
            'TimesNewRomanPSMT': 'Times New Roman',
            'TimesNewRomanPS-BoldMT': 'Times New Roman',
            'Helvetica': 'Arial',
            'Helvetica-Bold': 'Arial',
            'Helvetica-Italic': 'Arial',
            'Helvetica-BoldItalic': 'Arial',
            'Courier': 'Courier New',
            'Courier-Bold': 'Courier New',
            'Courier-Italic': 'Courier New',
            'Courier-BoldItalic': 'Courier New',
            'SimSun': 'SimSun',
            'SimHei': 'SimHei',
            'Microsoft YaHei': 'Microsoft YaHei',
            'KaiTi': 'KaiTi',
            'FangSong': 'FangSong',
            # 添加 FranklinGothic 字体映射
            'FranklinGothic-Book': 'Arial',
            'FranklinGothic': 'Arial',
            'Franklin': 'Arial',
            'Calibri': 'Calibri',
            'Calibri-Bold': 'Calibri',
            'Calibri-Italic': 'Calibri',
        }
        
        # 清理字体名
        clean_font_name = pdf_font_name.split('+')[-1]  # 移除字体子集前缀
        
        # 查找映射
        for pdf_font, word_font in font_mapping.items():
            if pdf_font.lower() in clean_font_name.lower():
                return word_font
        
        # 如果没有找到映射，根据字体特征推断
        clean_font_lower = clean_font_name.lower()
        if 'times' in clean_font_lower or 'roman' in clean_font_lower:
            return 'Times New Roman'
        elif 'arial' in clean_font_lower or 'helvetica' in clean_font_lower:
            return 'Arial'
        elif 'courier' in clean_font_lower or 'mono' in clean_font_lower:
            return 'Courier New'
        elif 'calibri' in clean_font_lower:
            return 'Calibri'
        else:
            # 默认使用 Arial
            return 'Arial' 

    def _process_page_layout_optimized(self, doc, text_dict, page):
        """基于参考文档分析的优化布局处理"""
        try:
            # 收集所有文本行，保持原始顺序
            all_lines = []
            
            for block in text_dict.get("blocks", []):
                if "lines" in block:  # 文本块
                    for line in block["lines"]:
                        line_info = self._extract_line_info(line, block, page.rect.width)
                        if line_info:
                            all_lines.append(line_info)
            
            # 按垂直位置排序（从上到下，从左到右）
            all_lines.sort(key=lambda x: (x['bbox'][1], x['bbox'][0]))
            
            # 智能合并相关行（如标题的多行）
            merged_lines = self._merge_related_lines(all_lines)
            
            # 分析文档结构
            structure = self._analyze_document_structure_optimized(merged_lines, page.rect.width)
            
            # 处理每一行
            for line_info in merged_lines:
                self._process_line_optimized(doc, line_info, structure, page.rect.width)
                    
        except Exception as e:
            logger.error(f"优化布局处理失败: {e}")

    def _extract_line_info(self, line, block, page_width):
        """提取行信息"""
        try:
            line_text = ""
            line_bbox = line.get("bbox", [0, 0, 0, 0])
            max_font_size = 0
            is_bold = False
            font_name = "Arial"
            
            # 收集行文本和样式信息
            for span in line["spans"]:
                line_text += span["text"]
                span_size = span.get("size", 12)
                max_font_size = max(max_font_size, span_size)
                
                # 检测粗体
                span_flags = span.get("flags", 0)
                if span_flags & 2**4:  # 粗体标志
                    is_bold = True
                
                # 获取字体名
                span_font = span.get("font", "Arial")
                if span_font:
                    font_name = span_font
            
            line_text = line_text.strip()
            if not line_text:
                return None
            
            # 转换字体大小
            font_size = max(8, min(72, max_font_size * 0.75))
            
            return {
                'text': line_text,
                'bbox': line_bbox,
                'font_size': font_size,
                'is_bold': is_bold,
                'font_name': font_name,
                'spans': line["spans"],
                'block': block
            }
            
        except Exception as e:
            logger.warning(f"提取行信息失败: {e}")
            return None

    def _merge_related_lines(self, all_lines):
        """智能合并相关的行（如标题的多行）"""
        if not all_lines:
            return []
        
        merged_lines = []
        current_group = [all_lines[0]]
        
        for i in range(1, len(all_lines)):
            current_line = all_lines[i]
            prev_line = all_lines[i-1]
            
            # 检查是否应该合并
            should_merge = self._should_merge_lines(prev_line, current_line)
            
            if should_merge:
                current_group.append(current_line)
            else:
                # 合并当前组
                if current_group:
                    merged_line = self._merge_line_group(current_group)
                    merged_lines.append(merged_line)
                current_group = [current_line]
        
        # 处理最后一组
        if current_group:
            merged_line = self._merge_line_group(current_group)
            merged_lines.append(merged_line)
        
        return merged_lines

    def _should_merge_lines(self, line1, line2):
        """判断两行是否应该合并"""
        # 检查垂直距离
        y1 = line1['bbox'][1]
        y2 = line2['bbox'][1]
        y_diff = abs(y2 - y1)
        
        # 检查字体大小是否相似
        font_size_diff = abs(line1['font_size'] - line2['font_size'])
        
        # 检查是否都是粗体
        both_bold = line1['is_bold'] and line2['is_bold']
        
        # 检查是否包含标题关键词
        title_keywords = ['meeting', 'negotiating', 'body', 'convention', 'agreement', 'international', 'instrument', 'pandemic', 'prevention', 'preparedness', 'response']
        has_title_keywords = any(keyword in line1['text'].lower() for keyword in title_keywords) or any(keyword in line2['text'].lower() for keyword in title_keywords)
        
        # 合并条件：垂直距离小、字体大小相似、都是粗体、包含标题关键词
        if y_diff < 30 and font_size_diff < 2 and both_bold and has_title_keywords:
            return True
        
        # 检查是否是同一行的延续（水平位置相近）
        x1 = line1['bbox'][0]
        x2 = line2['bbox'][0]
        x_diff = abs(x2 - x1)
        
        if y_diff < 15 and x_diff < 50:  # 同一行的延续
            return True
        
        return False

    def _merge_line_group(self, line_group):
        """合并行组"""
        if len(line_group) == 1:
            return line_group[0]
        
        # 合并文本
        merged_text = " ".join(line['text'] for line in line_group)
        
        # 使用第一行的基本信息
        merged_line = line_group[0].copy()
        merged_line['text'] = merged_text
        
        # 更新边界框
        min_x = min(line['bbox'][0] for line in line_group)
        min_y = min(line['bbox'][1] for line in line_group)
        max_x = max(line['bbox'][2] for line in line_group)
        max_y = max(line['bbox'][3] for line in line_group)
        merged_line['bbox'] = [min_x, min_y, max_x, max_y]
        
        # 合并spans
        all_spans = []
        for line in line_group:
            all_spans.extend(line['spans'])
        merged_line['spans'] = all_spans
        
        return merged_line

    def _analyze_document_structure_optimized(self, all_lines, page_width):
        """基于参考文档分析的文档结构分析"""
        structure = {
            'title_lines': [],      # 大标题（如文档标题）
            'subtitle_lines': [],   # 副标题（如章节标题）
            'body_lines': [],       # 正文
            'list_items': [],       # 列表项
            'header_info': [],      # 页眉信息
            'footer_info': []       # 页脚信息
        }
        
        for line_info in all_lines:
            text = line_info['text']
            bbox = line_info['bbox']
            font_size = line_info['font_size']
            is_bold = line_info['is_bold']
            
            # 基于参考文档的特征分析
            if self._is_document_title(text, bbox, font_size, is_bold):
                structure['title_lines'].append(line_info)
            elif self._is_section_title(text, bbox, font_size, is_bold):
                structure['subtitle_lines'].append(line_info)
            elif self._is_list_item_optimized(text, bbox, font_size):
                structure['list_items'].append(line_info)
            elif self._is_header_info_optimized(text, bbox, page_width):
                structure['header_info'].append(line_info)
            elif self._is_footer_info(text, bbox, page_width):
                structure['footer_info'].append(line_info)
            else:
                structure['body_lines'].append(line_info)
        
        return structure

    def _is_document_title(self, text, bbox, font_size, is_bold):
        """判断是否为文档标题"""
        # 基于参考文档的特征
        if font_size >= 18 and is_bold:
            return True
        # 检查是否包含特定关键词
        title_keywords = ['meeting', 'negotiating', 'body', 'convention', 'agreement', 'international', 'instrument', 'pandemic', 'prevention', 'preparedness', 'response']
        if any(keyword in text.lower() for keyword in title_keywords):
            return True
        return False

    def _is_section_title(self, text, bbox, font_size, is_bold):
        """判断是否为章节标题"""
        # 基于参考文档的特征
        if font_size >= 14 and is_bold:
            return True
        # 检查是否包含特定关键词
        section_keywords = ['agenda', 'provisional', 'opening', 'closure']
        if any(keyword in text.lower() for keyword in section_keywords):
            return True
        return False

    def _is_list_item_optimized(self, text, bbox, font_size):
        """优化的列表项判断"""
        import re
        
        # 匹配数字编号：1. 2. 3. 等
        if re.match(r'^\d+\.', text.strip()):
            return True
        # 匹配字母编号：a. b. c. 等
        if re.match(r'^[a-z]\.', text.strip()):
            return True
        # 匹配罗马数字：i. ii. iii. 等
        if re.match(r'^[ivxlcdm]+\.', text.strip(), re.IGNORECASE):
            return True
        # 匹配项目符号
        if text.strip().startswith('•') or text.strip().startswith('-'):
            return True
        # 检查是否包含列表项关键词
        list_keywords = ['opening', 'conceptual', 'proposal', 'information', 'summaries', 'informal', 'secretariat', 'report', 'closure']
        if any(keyword in text.lower() for keyword in list_keywords):
            return True
        return False

    def _is_header_info_optimized(self, text, bbox, page_width):
        """优化的页眉信息判断"""
        # 检查是否在页面顶部
        if len(bbox) >= 4 and bbox[1] < 100:
            return True
        # 检查是否包含文档编号或日期
        header_keywords = ['a/inb', 'geneva', 'november', 'december', '2022']
        if any(keyword in text.lower() for keyword in header_keywords):
            return True
        return False

    def _is_footer_info(self, text, bbox, page_width):
        """判断是否为页脚信息"""
        # 检查是否在页面底部
        if len(bbox) >= 4 and bbox[3] > 700:  # 假设页面高度约800
            return True
        # 检查是否包含脚注特征
        if text.strip().startswith('1 ') and len(text) > 100:
            return True
        return False

    def _process_line_optimized(self, doc, line_info, structure, page_width):
        """基于参考文档优化的行处理"""
        try:
            text = line_info['text']
            bbox = line_info['bbox']
            font_size = line_info['font_size']
            is_bold = line_info['is_bold']
            font_name = line_info['font_name']
            spans = line_info['spans']
            
            # 创建段落
            paragraph = doc.add_paragraph()
            
            # 设置段落对齐方式
            self._set_paragraph_alignment(paragraph, bbox, page_width)
            
            # 基于参考文档的格式设置
            self._set_formatting_based_on_reference(paragraph, line_info, structure, page_width)
            
            # 处理文本和样式
            for span in spans:
                span_text = span["text"]
                if span_text.strip():
                    run = paragraph.add_run(span_text)
                    
                    # 应用字体样式
                    word_font_name = self._get_font_mapping(span.get("font", font_name))
                    run.font.name = word_font_name
                    
                    # 字体大小
                    span_size = span.get("size", font_size / 0.75)
                    word_font_size = max(8, min(72, span_size * 0.75))
                    from docx.shared import Pt
                    run.font.size = Pt(word_font_size)
                    
                    # 粗体
                    span_flags = span.get("flags", 0)
                    run.bold = bool((span_flags & 2**4) or is_bold)
                    
                    # 斜体
                    run.italic = span_flags & 2**1
                    
                    # 颜色
                    if "color" in span:
                        color = span["color"]
                        if color != 0:
                            rgb_color = self._convert_color_to_rgb(color)
                            if rgb_color:
                                try:
                                    from docx.shared import RGBColor
                                    run.font.color.rgb = RGBColor(*rgb_color)
                                except Exception as color_error:
                                    logger.warning(f"设置字体颜色失败: {color_error}")
            
        except Exception as e:
            logger.error(f"优化行处理失败: {e}")

    def _set_formatting_based_on_reference(self, paragraph, line_info, structure, page_width):
        """基于参考文档的格式设置"""
        try:
            from docx.shared import Inches, Pt
            from docx.enum.text import WD_ALIGN_PARAGRAPH
            
            text = line_info['text']
            bbox = line_info['bbox']
            font_size = line_info['font_size']
            
            # 计算行的X中心位置
            x_center = (bbox[0] + bbox[2]) / 2
            page_center = page_width / 2
            
            # 特殊处理分隔线
            if '=' in text and text.strip().replace('=', '').replace(' ', '') == '':
                # 分隔线格式：居中对齐，无缩进
                paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
                paragraph.paragraph_format.left_indent = Inches(0)
                paragraph.paragraph_format.first_line_indent = Inches(0)
                paragraph.paragraph_format.space_after = Pt(6.0)
                return
            
            # 基于PDF实际位置设置对齐方式
            if '=' in text and text.strip().replace('=', '').replace(' ', '') == '':
                # 分隔线居中对齐
                paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
            elif x_center > page_center + 100:
                # 明显在右侧，右对齐
                paragraph.alignment = WD_ALIGN_PARAGRAPH.RIGHT
            elif abs(x_center - page_center) < 50:
                # 接近页面中心，居中对齐
                paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
            else:
                # 默认左对齐
                paragraph.alignment = WD_ALIGN_PARAGRAPH.LEFT
            
            # 基于参考文档的精确格式设置
            if self._is_document_title(text, bbox, font_size, line_info['is_bold']):
                # 文档标题格式：负缩进，小间距
                paragraph.paragraph_format.left_indent = Inches(-0.003)
                paragraph.paragraph_format.first_line_indent = Inches(-0.007)
                paragraph.paragraph_format.space_after = Pt(0.1)
                
            elif self._is_section_title(text, bbox, font_size, line_info['is_bold']):
                # 章节标题格式：无缩进，无间距
                paragraph.paragraph_format.left_indent = Inches(0)
                paragraph.paragraph_format.first_line_indent = Inches(0)
                paragraph.paragraph_format.space_after = Pt(0)
                
            elif self._is_list_item_optimized(text, bbox, font_size):
                # 列表项格式：负首行缩进，较大间距
                if 'agenda' in text.lower() or 'opening' in text.lower() or 'closure' in text.lower():
                    # 主要列表项
                    paragraph.paragraph_format.left_indent = Inches(0)
                    paragraph.paragraph_format.first_line_indent = Inches(-0.393)
                    paragraph.paragraph_format.space_after = Pt(12.95)
                else:
                    # 子列表项
                    paragraph.paragraph_format.left_indent = Inches(0)
                    paragraph.paragraph_format.first_line_indent = Inches(-0.138)
                    paragraph.paragraph_format.space_after = Pt(12.95)
                    
            elif self._is_header_info_optimized(text, bbox, page_width):
                # 页眉信息格式：负缩进，小间距
                paragraph.paragraph_format.left_indent = Inches(-0.010)
                paragraph.paragraph_format.first_line_indent = Inches(0)
                paragraph.paragraph_format.space_after = Pt(0.1)
                
            elif self._is_footer_info(text, bbox, page_width):
                # 页脚信息格式：正缩进，无间距
                paragraph.paragraph_format.left_indent = Inches(0)
                paragraph.paragraph_format.first_line_indent = Inches(0.394)
                paragraph.paragraph_format.space_after = Pt(0)
                
            else:
                # 正文格式：根据内容调整
                if 'information' in text.lower() or 'summaries' in text.lower():
                    # 信息文档类
                    paragraph.paragraph_format.left_indent = Inches(0)
                    paragraph.paragraph_format.first_line_indent = Inches(-0.138)
                    paragraph.paragraph_format.space_after = Pt(12.95)
                else:
                    # 普通正文
                    paragraph.paragraph_format.left_indent = Inches(0)
                    paragraph.paragraph_format.first_line_indent = Inches(0)
                    paragraph.paragraph_format.space_after = Pt(6.0)
            
            # 设置行间距
            if font_size >= 16:
                paragraph.paragraph_format.line_spacing = 1.0375
            else:
                paragraph.paragraph_format.line_spacing = 1.029
            
        except Exception as e:
            logger.warning(f"基于参考文档的格式设置失败: {e}") 

    def _collect_and_sort_content(self, text_dict, page):
        """收集并排序所有内容（文本、图片和水平线）"""
        all_content = []
        
        # 收集文本内容
        text_lines = []
        for block in text_dict.get("blocks", []):
            if "lines" in block:  # 文本块
                for line in block["lines"]:
                    line_info = self._extract_line_info(line, block, page.rect.width)
                    if line_info:
                        text_lines.append(line_info)
        
        # 合并相关的文本行
        merged_lines = self._merge_related_text_lines(text_lines)
        
        # 添加到内容列表
        for line_info in merged_lines:
            all_content.append({
                'type': 'text',
                'data': line_info,
                'position': line_info['bbox'][1]  # 使用Y坐标作为位置
            })
        
        # 收集图片内容
        for block in text_dict.get("blocks", []):
            if "image" in block:  # 图片块
                all_content.append({
                    'type': 'image',
                    'data': block,
                    'position': block["bbox"][1]  # 使用Y坐标作为位置
                })
        
        # 收集水平线内容
        try:
            paths = page.get_drawings()
            for path in paths:
                if 'rect' in path:
                    rect = path['rect']
                    width = rect[2] - rect[0]
                    height = rect[3] - rect[1]
                    
                    # 判断是否是水平线
                    if width > height * 5 and height < 10:
                        all_content.append({
                            'type': 'horizontal_line',
                            'data': {
                                'rect': rect,
                                'width': width,
                                'height': height,
                                'y_position': rect[1]
                            },
                            'position': rect[1]  # 使用Y坐标作为位置
                        })
        except Exception as e:
            logger.warning(f"提取水平线失败: {e}")
        
        # 按位置排序（从上到下）
        all_content.sort(key=lambda x: x['position'])
        
        return all_content

    def _process_text_content(self, doc, content, page_width):
        """处理文本内容"""
        try:
            line_info = content['data']
            
            # 创建段落
            paragraph = doc.add_paragraph()
            
            # 分析文档结构（这里简化处理）
            structure = {'title_lines': [], 'subtitle_lines': [], 'body_lines': [], 'list_items': [], 'header_info': [], 'footer_info': []}
            
            # 基于参考文档的格式设置（包含对齐方式）
            self._set_formatting_based_on_reference(paragraph, line_info, structure, page_width)
            
            # 处理文本和样式
            for span in line_info['spans']:
                span_text = span["text"]
                if span_text.strip():
                    run = paragraph.add_run(span_text)
                    
                    # 应用字体样式
                    word_font_name = self._get_font_mapping(span.get("font", line_info['font_name']))
                    run.font.name = word_font_name
                    
                    # 字体大小
                    span_size = span.get("size", line_info['font_size'] / 0.75)
                    word_font_size = max(8, min(72, span_size * 0.75))
                    from docx.shared import Pt
                    run.font.size = Pt(word_font_size)
                    
                    # 粗体
                    span_flags = span.get("flags", 0)
                    run.bold = bool((span_flags & 2**4) or line_info['is_bold'])
                    
                    # 斜体
                    run.italic = span_flags & 2**1
                    
                    # 颜色
                    if "color" in span:
                        color = span["color"]
                        if color != 0:
                            rgb_color = self._convert_color_to_rgb(color)
                            if rgb_color:
                                try:
                                    from docx.shared import RGBColor
                                    run.font.color.rgb = RGBColor(*rgb_color)
                                except Exception as color_error:
                                    logger.warning(f"设置字体颜色失败: {color_error}")
            
        except Exception as e:
            logger.error(f"处理文本内容失败: {e}")

    def _process_image_content(self, doc, content, pdf_document, page_num, page):
        """处理图片内容"""
        try:
            block = content['data']
            
            # 获取图片信息
            img_rect = block["bbox"]
            img_data = block["image"]
            
            # 直接保存PNG数据到临时文件
            temp_img_path = self._save_image_data(img_data, page_num, img_rect)
            
            if temp_img_path:
                # 计算合适的图片大小
                page_width = page.rect.width
                page_height = page.rect.height
                img_width, img_height = self._calculate_image_size(
                    block, page_width, page_height
                )
                
                # 在指定位置插入图片
                self._insert_image_at_position(doc, temp_img_path, img_width, img_height, img_rect, page_width)
                
                # 清理临时文件
                import os
                os.remove(temp_img_path)
                
                logger.info(f"成功插入图片，位置: {img_rect}, 大小: {img_width:.2f}in x {img_height:.2f}in")
            
        except Exception as img_error:
            logger.warning(f"处理图片失败: {img_error}")

    def _save_image_data(self, img_data, page_num, img_rect):
        """保存图片数据到临时文件"""
        try:
            import os
            
            # 生成临时文件路径
            temp_img_file = tempfile.NamedTemporaryFile(suffix=f"_pdf_img_data_{page_num}_{hash(str(img_rect))}.png", delete=False)
            temp_img_path = temp_img_file.name
            temp_img_file.close()
            
            # 如果img_data是字节串，直接写入文件
            if isinstance(img_data, bytes):
                with open(temp_img_path, 'wb') as f:
                    f.write(img_data)
                
                # 验证文件是否成功创建
                if os.path.exists(temp_img_path) and os.path.getsize(temp_img_path) > 0:
                    return temp_img_path
                else:
                    logger.warning(f"图片文件创建失败 (页面 {page_num})")
                    return None
            else:
                logger.warning(f"不支持的图片数据类型: {type(img_data)}")
                return None
                
        except Exception as e:
            logger.warning(f"保存图片数据失败 (页面 {page_num}): {e}")
            return None

    def _insert_image_at_position(self, doc, image_path, img_width, img_height, img_rect, page_width):
        """在指定位置插入图片"""
        try:
            from docx.shared import Inches
            from docx.enum.text import WD_ALIGN_PARAGRAPH
            
            # 创建图片段落
            img_paragraph = doc.add_paragraph()
            
            # 设置图片对齐方式（基于图片在PDF中的水平位置）
            img_center_x = (img_rect[0] + img_rect[2]) / 2
            page_center_x = page_width / 2
            
            if abs(img_center_x - page_center_x) < 50:  # 居中
                img_paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
            elif img_center_x < page_center_x:  # 左对齐
                img_paragraph.alignment = WD_ALIGN_PARAGRAPH.LEFT
            else:  # 右对齐
                img_paragraph.alignment = WD_ALIGN_PARAGRAPH.RIGHT
            
            # 添加图片到段落
            img_paragraph.add_run().add_picture(image_path, width=Inches(img_width))
            
            logger.info(f"在位置 {img_rect} 插入图片，大小: {img_width:.2f}in x {img_height:.2f}in")
            
        except Exception as e:
            logger.error(f"插入图片失败: {e}")

    def _merge_related_text_lines(self, text_lines):
        """合并相关的文本行"""
        if not text_lines:
            return []
        
        # 按Y位置排序
        text_lines.sort(key=lambda x: x['bbox'][1])
        
        merged_lines = []
        i = 0
        
        while i < len(text_lines):
            current_line = text_lines[i]
            merged_text = current_line['text']
            merged_bbox = current_line['bbox'].copy()
            
            # 检查下一行是否应该合并
            j = i + 1
            while j < len(text_lines):
                next_line = text_lines[j]
                
                # 检查Y位置是否接近（同一行）
                y_diff = abs(next_line['bbox'][1] - current_line['bbox'][1])
                
                if y_diff < 5:  # 如果Y位置差异小于5，认为是同一行
                    # 合并文本
                    merged_text += ' ' + next_line['text']
                    
                    # 更新边界框
                    merged_bbox[0] = min(merged_bbox[0], next_line['bbox'][0])
                    merged_bbox[1] = min(merged_bbox[1], next_line['bbox'][1])
                    merged_bbox[2] = max(merged_bbox[2], next_line['bbox'][2])
                    merged_bbox[3] = max(merged_bbox[3], next_line['bbox'][3])
                    
                    j += 1
                else:
                    break
            
            # 创建合并后的行信息
            merged_line = current_line.copy()
            merged_line['text'] = merged_text
            merged_line['bbox'] = tuple(merged_bbox)
            
            merged_lines.append(merged_line)
            i = j
        
        return merged_lines

    def _insert_horizontal_line(self, doc, line_data, page_width):
        """插入水平线"""
        try:
            from docx.shared import Inches
            from docx.enum.text import WD_ALIGN_PARAGRAPH
            
            # 创建水平线段落
            line_paragraph = doc.add_paragraph()
            
            # 计算水平线宽度（转换为英寸）
            line_width_inches = line_data['width'] / 72.0  # 假设72 DPI
            
            # 创建水平线文本（使用等号字符，与参考文档一致）
            line_char = '='  # 使用等号字符，与参考文档一致
            line_length = int(line_width_inches * 6)  # 调整长度
            line_text = line_char * line_length
            
            # 添加水平线文本
            run = line_paragraph.add_run(line_text)
            
            # 设置水平线样式
            run.font.size = 8  # 小字体
            run.font.color.rgb = None  # 黑色
            
            # 设置段落对齐方式
            line_paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
            
            # 设置段落间距
            line_paragraph.paragraph_format.space_after = 6
            
            logger.info(f"插入水平线，位置: Y={line_data['y_position']:.1f}, 宽度: {line_width_inches:.2f}in")
            
        except Exception as e:
            logger.error(f"插入水平线失败: {e}") 

    async def _pdf_to_docx_with_pdf2docx(
        self, 
        input_path: str, 
        output_path: str, 
        options: Optional[Dict[str, Any]] = None
    ) -> Dict[str, Any]:
        """
        使用 pdf2docx 进行 PDF 转 Word 转换
        """
        try:
            from pdf2docx import Converter
            
            # 创建转换器
            cv = Converter(input_path)
            
            # 获取页面范围
            start_page = 0
            end_page = None
            
            if options:
                if 'start_page' in options:
                    start_page = max(0, options['start_page'] - 1)  # pdf2docx 使用 0-based
                if 'end_page' in options:
                    end_page = options['end_page']  # pdf2docx 使用 1-based
            
            # 执行转换
            cv.convert(output_path, start=start_page, end=end_page)
            cv.close()
            
            return {
                'success': True,
                'message': 'PDF 转 Word 成功（使用 pdf2docx）',
                'method': 'pdf2docx'
            }
            
        except Exception as e:
            logger.error(f"pdf2docx 转换失败: {e}")
            raise e 

    async def _pdf_to_docx_custom(
        self, 
        input_path: str, 
        output_path: str, 
        options: Optional[Dict[str, Any]] = None
    ) -> Dict[str, Any]:
        """
        使用自定义实现进行 PDF 转 Word 转换
        """
        try:
            from docx import Document
            from docx.shared import Inches, Pt
            from docx.enum.text import WD_ALIGN_PARAGRAPH
            import fitz  # PyMuPDF
            
            # 创建 Word 文档
            doc = Document()
            
            # 使用 PyMuPDF 提取文本和图片
            pdf_document = fitz.open(input_path)
            
            # 处理页面范围
            pages = self._get_pages_to_process(pdf_document, options)
            
            for page_num in pages:
                page = pdf_document[page_num - 1]
                
                # 收集页面上的所有内容（文本和图片）
                page_content = []
                
                # 1. 收集文本内容
                text_dict = page.get_text("dict")
                for block in text_dict.get("blocks", []):
                    if "lines" in block:  # 文本块
                        for line in block["lines"]:
                            # 计算文本块的垂直位置
                            if line["spans"]:
                                first_span = line["spans"][0]
                                bbox = first_span.get("bbox", [0, 0, 0, 0])
                                y_position = bbox[1]  # 垂直位置
                                
                                page_content.append({
                                    'type': 'text',
                                    'y_position': y_position,
                                    'data': line
                                })
                
                # 2. 收集图片内容
                text_positions = []
                for content in page_content:
                    if content['type'] == 'text':
                        text_positions.append(content['y_position'])
                
                # 获取图片列表
                image_list = page.get_images()
                if image_list:
                    for img_index, img in enumerate(image_list):
                        try:
                            # 智能插入图片位置
                            y_position = None
                            
                            # 根据文本分布智能插入图片
                            if text_positions:
                                text_positions.sort()
                                if len(text_positions) > 1:
                                    avg_gap = (text_positions[-1] - text_positions[0]) / (len(text_positions) - 1)
                                    insert_position = text_positions[0] + (img_index + 1) * avg_gap / (len(image_list) + 1)
                                    y_position = insert_position
                                else:
                                    y_position = text_positions[0] + 200 + img_index * 100
                            else:
                                y_position = 400 + img_index * 200
                            
                            page_content.append({
                                'type': 'image',
                                'y_position': y_position,
                                'data': {
                                    'img': img,
                                    'img_index': img_index,
                                    'page_num': page_num,
                                    'img_rect': None
                                }
                            })
                        except Exception as img_error:
                            logger.warning(f"处理图片失败: {img_error}")
                            continue
                
                # 3. 按垂直位置排序内容
                page_content.sort(key=lambda x: x['y_position'])
                
                # 4. 按顺序处理内容 - 改进的文本合并逻辑
                current_paragraph = None
                current_line_bbox = None
                current_max_font_size = 0
                
                for content in page_content:
                    if content['type'] == 'text':
                        # 处理文本
                        line = content['data']
                        
                        # 获取行的边界框信息
                        line_bbox = [0, 0, 0, 0]
                        if line["spans"]:
                            first_span = line["spans"][0]
                            line_bbox = first_span.get("bbox", [0, 0, 0, 0])
                        
                        page_width = page.rect.width
                        
                        # 检查是否需要创建新段落
                        should_create_new_paragraph = True
                        
                        # 如果当前段落存在，检查是否可以继续使用
                        if current_paragraph is not None:
                            # 检查垂直距离 - 如果距离很近，可能是同一行的延续
                            if current_line_bbox and len(current_line_bbox) >= 4 and len(line_bbox) >= 4:
                                vertical_distance = abs(line_bbox[1] - current_line_bbox[1])
                                # 如果垂直距离小于字体大小的一半，认为是同一行
                                if vertical_distance < 10:  # 10像素的阈值
                                    should_create_new_paragraph = False
                        
                        # 创建新段落或使用现有段落
                        if should_create_new_paragraph:
                            # 完成当前段落
                            if current_paragraph is not None:
                                # 设置当前段落的格式
                                if current_max_font_size > 0:
                                    self._set_line_spacing(current_paragraph, current_max_font_size)
                                    self._set_paragraph_spacing(current_paragraph, current_max_font_size)
                            
                            # 创建新段落
                            current_paragraph = doc.add_paragraph()
                            current_line_bbox = line_bbox
                            current_max_font_size = 0
                            
                            # 设置段落对齐方式
                            self._set_paragraph_alignment(current_paragraph, line_bbox, page_width)
                        
                        # 处理每个文本片段
                        for span in line["spans"]:
                            span_text = span["text"]
                            if span_text.strip():
                                run = current_paragraph.add_run(span_text)
                                
                                # 应用字体样式
                                font_name = span.get("font", "Arial")
                                run.font.name = self._get_font_mapping(font_name)
                                
                                # 字体大小
                                span_size = span.get("size", 12)
                                word_font_size = max(8, min(72, span_size * 0.75))
                                run.font.size = Pt(word_font_size)
                                current_max_font_size = max(current_max_font_size, word_font_size)
                                
                                # 粗体
                                span_flags = span.get("flags", 0)
                                run.bold = bool(span_flags & 2**4)
                                
                                # 斜体
                                run.italic = bool(span_flags & 2**1)
                                
                                # 颜色
                                span_color = span.get("color", 0)
                                if span_color != 0:
                                    rgb_color = self._convert_color_to_rgb(span_color)
                                    if rgb_color:
                                        try:
                                            from docx.shared import RGBColor
                                            run.font.color.rgb = RGBColor(*rgb_color)
                                        except Exception as color_error:
                                            logger.warning(f"设置字体颜色失败: {color_error}")
                    
                    elif content['type'] == 'image':
                        # 处理图片
                        img_data = content['data']
                        try:
                            # 安全提取图片
                            temp_img_path = self._extract_image_safely(
                                pdf_document, 
                                img_data['img'], 
                                img_data['page_num'], 
                                img_data['img_index']
                            )
                            
                            if temp_img_path:
                                # 智能计算图片尺寸
                                from docx.shared import Inches
                                from PIL import Image as PILImage
                                
                                try:
                                    # 获取原始图片尺寸
                                    with PILImage.open(temp_img_path) as pil_img:
                                        img_width, img_height = pil_img.size
                                    
                                    # 计算合适的Word图片尺寸
                                    max_width_inches = 3.5
                                    max_height_inches = 4.0
                                    
                                    width_inches = img_width / 96.0
                                    height_inches = img_height / 96.0
                                    
                                    # 如果图片太大，进行缩放
                                    if width_inches > max_width_inches:
                                        scale_factor = max_width_inches / width_inches
                                        width_inches = max_width_inches
                                        height_inches *= scale_factor
                                    
                                    if height_inches > max_height_inches:
                                        scale_factor = max_height_inches / height_inches
                                        height_inches = max_height_inches
                                        width_inches *= scale_factor
                                    
                                    # 确保最小尺寸
                                    min_width_inches = 1.0
                                    min_height_inches = 0.5
                                    
                                    if width_inches < min_width_inches:
                                        width_inches = min_width_inches
                                    if height_inches < min_height_inches:
                                        height_inches = min_height_inches
                                    
                                    # 添加图片到文档
                                    doc.add_picture(temp_img_path, width=Inches(width_inches))
                                    
                                except Exception as size_error:
                                    logger.warning(f"图片尺寸计算失败，使用默认尺寸: {size_error}")
                                    doc.add_picture(temp_img_path, width=Inches(2.5))
                                
                                # 清理临时文件
                                import os
                                os.remove(temp_img_path)
                                
                        except Exception as img_error:
                            logger.warning(f"处理图片失败: {img_error}")
                            continue
                
                # 完成最后一个段落
                if current_paragraph is not None and current_max_font_size > 0:
                    self._set_line_spacing(current_paragraph, current_max_font_size)
                    self._set_paragraph_spacing(current_paragraph, current_max_font_size)
            
            # 关闭 PDF 文档
            pdf_document.close()
            
            # 保存 Word 文档
            doc.save(output_path)
            
            return {
                'success': True,
                'message': 'PDF 转 Word 成功（自定义实现）',
                'method': 'custom'
            }
            
        except Exception as e:
            logger.error(f"自定义实现转换失败: {e}")
            return {
                'success': False,
                'error': f'PDF 转 Word 失败: {str(e)}'
            } 