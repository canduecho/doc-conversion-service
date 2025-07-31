"""
Markdown 转换器
支持 Markdown 与 PDF 之间的转换
"""

import os
import re
from pathlib import Path
from typing import Dict, Any, Optional
from loguru import logger

# 导入依赖
try:
    import markdown
    from weasyprint import HTML, CSS
    import fitz  # PyMuPDF
    from docx import Document
    from docx.shared import Pt
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from bs4 import BeautifulSoup
    from openpyxl import Workbook
    from openpyxl.styles import Font, Alignment, PatternFill, Border, Side
    from openpyxl.utils import get_column_letter
    from pptx import Presentation
    from pptx.util import Inches, Pt as PptPt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    DEPENDENCIES_AVAILABLE = True
except ImportError as e:
    logger.warning(f"Markdown 转换器依赖未安装: {e}")
    DEPENDENCIES_AVAILABLE = False


class MarkdownConverter:
    """
    Markdown 转换器
    支持 Markdown ↔ PDF 转换
    """
    
    def __init__(self):
        """初始化转换器"""
        if not DEPENDENCIES_AVAILABLE:
            logger.error("Markdown 转换器依赖未安装，无法使用")
            return
            
        self.default_css = """
        body {
            font-family: 'Segoe UI', Tahoma, Geneva, Verdana, sans-serif;
            line-height: 1.6;
            margin: 2em;
            color: #333;
        }
        h1, h2, h3, h4, h5, h6 {
            color: #2c3e50;
            margin-top: 1.5em;
            margin-bottom: 0.5em;
        }
        h1 { font-size: 2em; border-bottom: 2px solid #3498db; }
        h2 { font-size: 1.5em; border-bottom: 1px solid #bdc3c7; }
        h3 { font-size: 1.3em; }
        h4 { font-size: 1.1em; }
        code {
            background-color: #f8f9fa;
            padding: 0.2em 0.4em;
            border-radius: 3px;
            font-family: 'Courier New', monospace;
            font-size: 0.9em;
        }
        pre {
            background-color: #f8f9fa;
            padding: 1em;
            border-radius: 5px;
            overflow-x: auto;
            border-left: 4px solid #3498db;
        }
        pre code {
            background-color: transparent;
            padding: 0;
        }
        blockquote {
            border-left: 4px solid #bdc3c7;
            margin: 1em 0;
            padding-left: 1em;
            color: #7f8c8d;
        }
        table {
            border-collapse: collapse;
            width: 100%;
            margin: 1em 0;
        }
        th, td {
            border: 1px solid #ddd;
            padding: 8px;
            text-align: left;
        }
        th {
            background-color: #f2f2f2;
            font-weight: bold;
        }
        img {
            max-width: 100%;
            height: auto;
        }
        a {
            color: #3498db;
            text-decoration: none;
        }
        a:hover {
            text-decoration: underline;
        }
        ul, ol {
            padding-left: 2em;
        }
        li {
            margin: 0.5em 0;
        }
        hr {
            border: none;
            border-top: 1px solid #ddd;
            margin: 2em 0;
        }
        """
    
    async def markdown_to_docx(
        self,
        input_path: str,
        output_path: str,
        options: Optional[Dict[str, Any]] = None
    ) -> Dict[str, Any]:
        """
        将 Markdown 转换为 DOCX
        
        Args:
            input_path: 输入 Markdown 文件路径
            output_path: 输出 DOCX 文件路径
            options: 转换选项
            
        Returns:
            转换结果字典
        """
        try:
            if not DEPENDENCIES_AVAILABLE:
                return {
                    'success': False,
                    'error': 'Markdown 转换器依赖未安装，请安装 python-docx 和 beautifulsoup4'
                }
            
            logger.info(f"开始 Markdown 转 DOCX: {input_path}")
            
            # 读取 Markdown 文件
            with open(input_path, 'r', encoding='utf-8') as f:
                md_content = f.read()
            
            # 转换 Markdown 为 HTML
            html_content = markdown.markdown(
                md_content,
                extensions=['tables', 'fenced_code', 'codehilite', 'toc', 'nl2br']
            )
            
            # 创建 DOCX 文档
            doc = Document()
            
            # 解析 HTML 并转换为 DOCX
            soup = BeautifulSoup(html_content, 'html.parser')
            
            # 处理每个元素
            for element in soup.find_all(['h1', 'h2', 'h3', 'h4', 'h5', 'h6', 'p', 'ul', 'ol', 'blockquote', 'pre', 'hr', 'table']):
                if element.name in ['h1', 'h2', 'h3', 'h4', 'h5', 'h6']:
                    self._add_heading(doc, element)
                elif element.name == 'p':
                    self._add_paragraph(doc, element)
                elif element.name in ['ul', 'ol']:
                    self._add_list(doc, element)
                elif element.name == 'blockquote':
                    self._add_blockquote(doc, element)
                elif element.name == 'pre':
                    self._add_code_block(doc, element)
                elif element.name == 'hr':
                    self._add_horizontal_rule(doc)
                elif element.name == 'table':
                    self._add_table(doc, element)
            
            # 保存文档
            doc.save(output_path)
            
            logger.info(f"Markdown 转 DOCX 成功: {output_path}")
            return {
                'success': True,
                'output_path': output_path,
                'output_filename': Path(output_path).name
            }
            
        except Exception as e:
            logger.error(f"Markdown 转 DOCX 失败: {e}")
            return {
                'success': False,
                'error': f'Markdown 转 DOCX 失败: {str(e)}'
            }

    def _add_heading(self, doc, element):
        """添加标题到 DOCX 文档"""
        try:
            heading = doc.add_heading(element.get_text().strip(), level=int(element.name[1]))
            # 设置标题样式
            heading.alignment = WD_ALIGN_PARAGRAPH.LEFT
            return heading
        except Exception as e:
            logger.error(f"添加标题失败: {e}")
            # 如果失败，添加为普通段落
            return doc.add_paragraph(element.get_text().strip())

    def _add_paragraph(self, doc, element):
        """添加段落到 DOCX 文档"""
        try:
            paragraph = doc.add_paragraph(element.get_text().strip())
            paragraph.alignment = WD_ALIGN_PARAGRAPH.LEFT
            return paragraph
        except Exception as e:
            logger.error(f"添加段落失败: {e}")
            return None

    def _add_list(self, doc, element):
        """添加列表到 DOCX 文档"""
        try:
            items = []
            for li in element.find_all('li'):
                items.append(li.get_text().strip())
            
            if element.name == 'ol':  # 有序列表
                for i, item in enumerate(items, 1):
                    doc.add_paragraph(f"{i}. {item}", style='List Number')
            else:  # 无序列表
                for item in items:
                    doc.add_paragraph(f"• {item}", style='List Bullet')
        except Exception as e:
            logger.error(f"添加列表失败: {e}")

    def _add_blockquote(self, doc, element):
        """添加引用块到 DOCX 文档"""
        try:
            quote = doc.add_paragraph(element.get_text().strip())
            quote.alignment = WD_ALIGN_PARAGRAPH.LEFT
            # 可以添加引用样式（如果有的话）
            return quote
        except Exception as e:
            logger.error(f"添加引用失败: {e}")
            return None

    def _add_code_block(self, doc, element):
        """添加代码块到 DOCX 文档"""
        try:
            code_text = element.get_text().strip()
            code_paragraph = doc.add_paragraph(code_text)
            code_paragraph.alignment = WD_ALIGN_PARAGRAPH.LEFT
            # 可以设置代码样式（如果有的话）
            return code_paragraph
        except Exception as e:
            logger.error(f"添加代码块失败: {e}")
            return None

    def _add_horizontal_rule(self, doc):
        """添加水平分割线到 DOCX 文档"""
        try:
            # 添加一个空段落作为分割线
            hr_paragraph = doc.add_paragraph("─" * 50)
            hr_paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
            return hr_paragraph
        except Exception as e:
            logger.error(f"添加分割线失败: {e}")
            return None

    def _add_table(self, doc, element):
        """添加表格到 DOCX 文档"""
        try:
            rows = element.find_all('tr')
            if not rows:
                return None
            
            # 计算表格大小
            max_cols = 0
            for row in rows:
                cols = row.find_all(['td', 'th'])
                max_cols = max(max_cols, len(cols))
            
            if max_cols == 0:
                return None
            
            # 创建表格
            table = doc.add_table(rows=len(rows), cols=max_cols)
            table.style = 'Table Grid'
            
            # 填充表格数据
            for i, row in enumerate(rows):
                cells = row.find_all(['td', 'th'])
                for j, cell in enumerate(cells):
                    if j < max_cols:
                        table.cell(i, j).text = cell.get_text().strip()
            
            return table
        except Exception as e:
            logger.error(f"添加表格失败: {e}")
            return None

    async def markdown_to_pptx(
        self,
        input_path: str,
        output_path: str,
        options: Optional[Dict[str, Any]] = None
    ) -> Dict[str, Any]:
        """
        将 Markdown 转换为 PPTX
        
        Args:
            input_path: 输入 Markdown 文件路径
            output_path: 输出 PPTX 文件路径
            options: 转换选项
            
        Returns:
            转换结果字典
        """
        try:
            if not DEPENDENCIES_AVAILABLE:
                return {
                    'success': False,
                    'error': 'Markdown 转换器依赖未安装，请安装 python-pptx 和 beautifulsoup4'
                }
            
            logger.info(f"开始 Markdown 转 PPTX: {input_path}")
            
            # 读取 Markdown 文件
            with open(input_path, 'r', encoding='utf-8') as f:
                md_content = f.read()
            
            # 转换 Markdown 为 HTML
            html_content = markdown.markdown(
                md_content,
                extensions=['tables', 'fenced_code', 'codehilite', 'toc', 'nl2br']
            )
            
            # 创建 PPTX 演示文稿
            prs = Presentation()
            
            # 解析 HTML 并转换为 PPTX
            soup = BeautifulSoup(html_content, 'html.parser')
            
            # 设置默认样式
            title_font_size = PptPt(32)
            subtitle_font_size = PptPt(24)
            content_font_size = PptPt(18)
            body_font_size = PptPt(14)
            
            # 处理每个元素
            current_slide = None
            slide_content = []
            
            for element in soup.find_all(['h1', 'h2', 'h3', 'h4', 'h5', 'h6', 'p', 'ul', 'ol', 'blockquote', 'pre', 'hr', 'table']):
                if element.name in ['h1', 'h2']:
                    # H1 和 H2 创建新幻灯片
                    if current_slide:
                        # 处理当前幻灯片的内容
                        self._process_slide_content(current_slide, slide_content)
                        slide_content = []
                    
                    # 创建新幻灯片 - 尝试不同的布局
                    slide_layout = None
                    
                    # 首先尝试布局 0 (通常是标题页)
                    if len(prs.slide_layouts) > 0:
                        layout0 = prs.slide_layouts[0]
                        has_title_top = False
                        has_body = False
                        
                        for shape in layout0.placeholders:
                            if shape.placeholder_format.type == 1:  # TITLE
                                if shape.top < 2000000:  # 小于2英寸
                                    has_title_top = True
                            elif shape.placeholder_format.type == 2:  # BODY
                                has_body = True
                        
                        if has_title_top and has_body:
                            slide_layout = layout0
                            logger.info("使用布局 0 (标题在顶部)")
                    
                    # 如果布局 0 不合适，尝试布局 1
                    if not slide_layout and len(prs.slide_layouts) > 1:
                        layout1 = prs.slide_layouts[1]
                        has_title_top = False
                        has_body = False
                        
                        for shape in layout1.placeholders:
                            if shape.placeholder_format.type == 1:  # TITLE
                                if shape.top < 2000000:  # 小于2英寸
                                    has_title_top = True
                            elif shape.placeholder_format.type == 2:  # BODY
                                has_body = True
                        
                        if has_title_top and has_body:
                            slide_layout = layout1
                            logger.info("使用布局 1 (标题在顶部)")
                    
                    # 如果还是不合适，尝试布局 2
                    if not slide_layout and len(prs.slide_layouts) > 2:
                        layout2 = prs.slide_layouts[2]
                        has_title_top = False
                        has_body = False
                        
                        for shape in layout2.placeholders:
                            if shape.placeholder_format.type == 1:  # TITLE
                                if shape.top < 2000000:  # 小于2英寸
                                    has_title_top = True
                            elif shape.placeholder_format.type == 2:  # BODY
                                has_body = True
                        
                        if has_title_top and has_body:
                            slide_layout = layout2
                            logger.info("使用布局 2 (标题在顶部)")
                    
                    # 如果还是不合适，尝试布局 3
                    if not slide_layout and len(prs.slide_layouts) > 3:
                        layout3 = prs.slide_layouts[3]
                        has_title_top = False
                        has_body = False
                        
                        for shape in layout3.placeholders:
                            if shape.placeholder_format.type == 1:  # TITLE
                                if shape.top < 2000000:  # 小于2英寸
                                    has_title_top = True
                            elif shape.placeholder_format.type == 2:  # BODY
                                has_body = True
                        
                        if has_title_top and has_body:
                            slide_layout = layout3
                            logger.info("使用布局 3 (标题在顶部)")
                    
                    # 如果还是不合适，使用默认布局
                    if not slide_layout:
                        slide_layout = prs.slide_layouts[1]
                        logger.info("使用默认布局")
                    
                    current_slide = prs.slides.add_slide(slide_layout)
                    
                    # 设置标题
                    title = current_slide.shapes.title
                    if title:
                        title.text = element.get_text().strip()
                        
                        # 设置标题样式
                        title_frame = title.text_frame
                        if title_frame.paragraphs:
                            title_frame.paragraphs[0].font.size = title_font_size
                            title_frame.paragraphs[0].font.bold = True
                            title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                        else:
                            # 如果没有段落，创建一个
                            p = title_frame.add_paragraph()
                            p.text = element.get_text().strip()
                            p.font.size = title_font_size
                            p.font.bold = True
                            p.alignment = PP_ALIGN.CENTER
                    
                elif element.name in ['h3', 'h4', 'h5', 'h6']:
                    # H3-H6 添加到当前幻灯片作为子标题
                    if current_slide:
                        slide_content.append({
                            'type': 'subtitle',
                            'text': element.get_text().strip(),
                            'level': int(element.name[1])
                        })
                    else:
                        # 如果没有当前幻灯片，创建一个
                        slide_layout = prs.slide_layouts[1]
                        current_slide = prs.slides.add_slide(slide_layout)
                        slide_content.append({
                            'type': 'subtitle',
                            'text': element.get_text().strip(),
                            'level': int(element.name[1])
                        })
                
                elif element.name == 'p':
                    # 段落添加到当前幻灯片
                    if current_slide:
                        slide_content.append({
                            'type': 'paragraph',
                            'text': element.get_text().strip()
                        })
                    else:
                        # 如果没有当前幻灯片，创建一个
                        slide_layout = prs.slide_layouts[1]
                        current_slide = prs.slides.add_slide(slide_layout)
                        slide_content.append({
                            'type': 'paragraph',
                            'text': element.get_text().strip()
                        })
                
                elif element.name in ['ul', 'ol']:
                    # 列表添加到当前幻灯片
                    if current_slide:
                        slide_content.append({
                            'type': 'list',
                            'items': [li.get_text().strip() for li in element.find_all('li', recursive=False)],
                            'ordered': element.name == 'ol'
                        })
                    else:
                        # 如果没有当前幻灯片，创建一个
                        slide_layout = prs.slide_layouts[1]
                        current_slide = prs.slides.add_slide(slide_layout)
                        slide_content.append({
                            'type': 'list',
                            'items': [li.get_text().strip() for li in element.find_all('li', recursive=False)],
                            'ordered': element.name == 'ol'
                        })
                
                elif element.name == 'blockquote':
                    # 引用块添加到当前幻灯片
                    if current_slide:
                        slide_content.append({
                            'type': 'quote',
                            'text': element.get_text().strip()
                        })
                    else:
                        # 如果没有当前幻灯片，创建一个
                        slide_layout = prs.slide_layouts[1]
                        current_slide = prs.slides.add_slide(slide_layout)
                        slide_content.append({
                            'type': 'quote',
                            'text': element.get_text().strip()
                        })
                
                elif element.name == 'pre':
                    # 代码块添加到当前幻灯片
                    if current_slide:
                        slide_content.append({
                            'type': 'code',
                            'text': element.get_text().strip()
                        })
                    else:
                        # 如果没有当前幻灯片，创建一个
                        slide_layout = prs.slide_layouts[1]
                        current_slide = prs.slides.add_slide(slide_layout)
                        slide_content.append({
                            'type': 'code',
                            'text': element.get_text().strip()
                        })
                
                elif element.name == 'hr':
                    # 分割线创建新幻灯片
                    if current_slide:
                        # 处理当前幻灯片的内容
                        self._process_slide_content(current_slide, slide_content)
                        slide_content = []
                    
                    # 创建新幻灯片
                    slide_layout = prs.slide_layouts[1]
                    current_slide = prs.slides.add_slide(slide_layout)
                
                elif element.name == 'table':
                    # 表格创建新幻灯片
                    if current_slide:
                        # 处理当前幻灯片的内容
                        self._process_slide_content(current_slide, slide_content)
                        slide_content = []
                    
                    # 创建新幻灯片用于表格
                    slide_layout = prs.slide_layouts[1]
                    current_slide = prs.slides.add_slide(slide_layout)
                    
                    # 处理表格
                    self._add_pptx_table(current_slide, element)
            
            # 处理最后一个幻灯片的内容
            if current_slide and slide_content:
                self._process_slide_content(current_slide, slide_content)
            
            # 保存演示文稿
            prs.save(output_path)
            
            logger.info(f"Markdown 转 PPTX 成功: {output_path}")
            return {
                'success': True,
                'output_path': output_path,
                'output_filename': Path(output_path).name
            }
            
        except Exception as e:
            logger.error(f"Markdown 转 PPTX 失败: {e}")
            return {
                'success': False,
                'error': f'Markdown 转 PPTX 失败: {str(e)}'
            }

    def _process_slide_content(self, slide, content_items):
        """处理幻灯片内容"""
        try:
            # 获取内容占位符
            content_placeholder = None
            
            # 优先查找 BODY 占位符 (类型 2)
            for shape in slide.placeholders:
                if shape.placeholder_format.type == 2:  # BODY 占位符
                    content_placeholder = shape
                    logger.info(f"找到 BODY 占位符")
                    break
            
            # 其次查找 TEXT 占位符 (类型 3)
            if not content_placeholder:
                for shape in slide.placeholders:
                    if shape.placeholder_format.type == 3:  # TEXT 占位符
                        content_placeholder = shape
                        logger.info(f"找到 TEXT 占位符")
                        break
            
            # 再次查找 CONTENT 占位符 (类型 12)
            if not content_placeholder:
                for shape in slide.placeholders:
                    if shape.placeholder_format.type == 12:  # CONTENT 占位符
                        content_placeholder = shape
                        logger.info(f"找到 CONTENT 占位符")
                        break
            
            # 如果还是没找到，尝试查找其他内容相关占位符
            if not content_placeholder:
                for shape in slide.placeholders:
                    if shape.placeholder_format.type in [4, 5]:  # CHART, TABLE
                        content_placeholder = shape
                        logger.info(f"找到其他内容占位符，类型: {shape.placeholder_format.type}")
                        break
            
            # 如果仍然没找到，强制使用第一个非标题占位符
            if not content_placeholder:
                for shape in slide.placeholders:
                    if shape.placeholder_format.type != 1:  # 不是标题占位符
                        content_placeholder = shape
                        logger.info(f"强制使用占位符，类型: {shape.placeholder_format.type}")
                        break
            
            # 如果仍然没找到，创建文本框（这是最后的选择）
            if not content_placeholder:
                logger.warning("未找到内容占位符，创建文本框")
                left = Inches(1)
                top = Inches(2)
                width = Inches(8)
                height = Inches(5)
                textbox = slide.shapes.add_textbox(left, top, width, height)
                text_frame = textbox.text_frame
            else:
                text_frame = content_placeholder.text_frame
                logger.info(f"使用占位符，类型: {content_placeholder.placeholder_format.type}")
            
            # 清空现有内容
            text_frame.clear()
            
            # 添加内容
            for item in content_items:
                if item['type'] == 'subtitle':
                    self._add_pptx_subtitle(text_frame, item['text'], item['level'])
                elif item['type'] == 'paragraph':
                    self._add_pptx_paragraph(text_frame, item['text'])
                elif item['type'] == 'list':
                    self._add_pptx_list(text_frame, item['items'], item['ordered'])
                elif item['type'] == 'quote':
                    self._add_pptx_quote(text_frame, item['text'])
                elif item['type'] == 'code':
                    self._add_pptx_code(text_frame, item['text'])
                    
            logger.info(f"成功添加 {len(content_items)} 个内容项到占位符")
                    
        except Exception as e:
            logger.error(f"处理幻灯片内容失败: {e}")
            raise

    def _add_pptx_subtitle(self, text_frame, text, level):
        """添加 PPTX 子标题"""
        if text:
            p = text_frame.add_paragraph()
            p.text = text
            p.font.size = PptPt(20 - level * 2)
            p.font.bold = True
            p.font.color.rgb = RGBColor(0, 0, 139)  # 深蓝色

    def _add_pptx_paragraph(self, text_frame, text):
        """添加 PPTX 段落"""
        if text:
            p = text_frame.add_paragraph()
            p.text = text
            p.font.size = PptPt(14)
            p.font.name = 'Calibri'

    def _add_pptx_list(self, text_frame, items, ordered):
        """添加 PPTX 列表"""
        for i, item in enumerate(items):
            if item:
                p = text_frame.add_paragraph()
                if ordered:
                    p.text = f"{i + 1}. {item}"
                else:
                    p.text = f"• {item}"
                p.font.size = PptPt(14)
                p.font.name = 'Calibri'
                p.level = 1  # 缩进级别

    def _add_pptx_quote(self, text_frame, text):
        """添加 PPTX 引用"""
        if text:
            p = text_frame.add_paragraph()
            p.text = f"> {text}"
            p.font.size = PptPt(12)
            p.font.italic = True
            p.font.color.rgb = RGBColor(128, 128, 128)  # 灰色

    def _add_pptx_code(self, text_frame, text):
        """添加 PPTX 代码块"""
        if text:
            p = text_frame.add_paragraph()
            p.text = text
            p.font.size = PptPt(10)
            p.font.name = 'Courier New'
            p.font.color.rgb = RGBColor(0, 0, 0)  # 黑色

    def _add_pptx_table(self, slide, element):
        """添加 PPTX 表格"""
        try:
            rows = element.find_all('tr')
            if rows:
                # 计算表格大小
                max_cols = max(len(row.find_all(['td', 'th'])) for row in rows)
                
                # 创建表格
                table = slide.shapes.add_table(
                    rows=len(rows), 
                    cols=max_cols,
                    left=Inches(1),
                    top=Inches(2),
                    width=Inches(8),
                    height=Inches(4)
                ).table
                
                # 填充表格数据
                for i, row in enumerate(rows):
                    cells = row.find_all(['td', 'th'])
                    for j, cell in enumerate(cells):
                        if j < max_cols:
                            table.cell(i, j).text = cell.get_text().strip()
                            
                            # 设置表头样式
                            if cell.name == 'th':
                                table.cell(i, j).text_frame.paragraphs[0].font.bold = True
                                table.cell(i, j).text_frame.paragraphs[0].font.size = PptPt(12)
                            else:
                                table.cell(i, j).text_frame.paragraphs[0].font.size = PptPt(10)
                                
        except Exception as e:
            logger.error(f"添加 PPTX 表格失败: {e}")

    async def markdown_to_xlsx(
        self,
        input_path: str,
        output_path: str,
        options: Optional[Dict[str, Any]] = None
    ) -> Dict[str, Any]:
        """
        将 Markdown 转换为 XLSX
        
        Args:
            input_path: 输入 Markdown 文件路径
            output_path: 输出 XLSX 文件路径
            options: 转换选项
            
        Returns:
            转换结果字典
        """
        try:
            if not DEPENDENCIES_AVAILABLE:
                return {
                    'success': False,
                    'error': 'Markdown 转换器依赖未安装，请安装 openpyxl 和 beautifulsoup4'
                }
            
            logger.info(f"开始 Markdown 转 XLSX: {input_path}")
            
            # 读取 Markdown 文件
            with open(input_path, 'r', encoding='utf-8') as f:
                md_content = f.read()
            
            # 转换 Markdown 为 HTML
            html_content = markdown.markdown(
                md_content,
                extensions=['tables', 'fenced_code', 'codehilite', 'toc', 'nl2br']
            )
            
            # 创建 Excel 工作簿
            wb = Workbook()
            ws = wb.active
            ws.title = "Markdown Content"
            
            # 解析 HTML 并转换为 XLSX
            soup = BeautifulSoup(html_content, 'html.parser')
            
            # 设置样式
            header_font = Font(bold=True, size=12, color="FFFFFF")
            header_fill = PatternFill(start_color="366092", end_color="366092", fill_type="solid")
            content_font = Font(size=11)
            border = Border(
                left=Side(style='thin'),
                right=Side(style='thin'),
                top=Side(style='thin'),
                bottom=Side(style='thin')
            )
            
            current_row = 1
            
            # 处理每个元素
            for element in soup.find_all(['h1', 'h2', 'h3', 'h4', 'h5', 'h6', 'p', 'ul', 'ol', 'blockquote', 'pre', 'hr', 'table']):
                if element.name in ['h1', 'h2', 'h3', 'h4', 'h5', 'h6']:
                    current_row = self._add_excel_heading(ws, element, current_row, header_font, header_fill, border)
                elif element.name == 'p':
                    current_row = self._add_excel_paragraph(ws, element, current_row, content_font, border)
                elif element.name in ['ul', 'ol']:
                    current_row = self._add_excel_list(ws, element, current_row, content_font, border)
                elif element.name == 'blockquote':
                    current_row = self._add_excel_blockquote(ws, element, current_row, content_font, border)
                elif element.name == 'pre':
                    current_row = self._add_excel_code_block(ws, element, current_row, content_font, border)
                elif element.name == 'hr':
                    current_row = self._add_excel_horizontal_rule(ws, current_row, border)
                elif element.name == 'table':
                    current_row = self._add_excel_table(ws, element, current_row, header_font, header_fill, content_font, border)
            
            # 自动调整列宽
            for column in ws.columns:
                max_length = 0
                column_letter = get_column_letter(column[0].column)
                for cell in column:
                    try:
                        if len(str(cell.value)) > max_length:
                            max_length = len(str(cell.value))
                    except:
                        pass
                adjusted_width = min(max_length + 2, 50)
                ws.column_dimensions[column_letter].width = adjusted_width
            
            # 保存文档
            wb.save(output_path)
            
            logger.info(f"Markdown 转 XLSX 成功: {output_path}")
            return {
                'success': True,
                'output_path': output_path,
                'output_filename': Path(output_path).name
            }
            
        except Exception as e:
            logger.error(f"Markdown 转 XLSX 失败: {e}")
            return {
                'success': False,
                'error': f'Markdown 转 XLSX 失败: {str(e)}'
            }

    def _add_excel_heading(self, ws, element, current_row, header_font, header_fill, border):
        """添加 Excel 标题"""
        level = int(element.name[1])
        text = element.get_text().strip()
        if text:
            # 设置标题单元格
            cell = ws.cell(row=current_row, column=1, value=text)
            cell.font = header_font
            cell.fill = header_fill
            cell.border = border
            cell.alignment = Alignment(horizontal='left', vertical='center')
            
            # 合并单元格（标题占整行）
            ws.merge_cells(f'A{current_row}:Z{current_row}')
            
            return current_row + 1
        return current_row

    def _add_excel_paragraph(self, ws, element, current_row, content_font, border):
        """添加 Excel 段落"""
        text = element.get_text().strip()
        if text:
            cell = ws.cell(row=current_row, column=1, value=text)
            cell.font = content_font
            cell.border = border
            cell.alignment = Alignment(horizontal='left', vertical='top', wrap_text=True)
            return current_row + 1
        return current_row

    def _add_excel_list(self, ws, element, current_row, content_font, border):
        """添加 Excel 列表"""
        is_ordered = element.name == 'ol'
        
        for li in element.find_all('li', recursive=False):
            text = li.get_text().strip()
            if text:
                # 添加列表标记
                marker = f"{current_row}. " if is_ordered else "• "
                cell = ws.cell(row=current_row, column=1, value=marker + text)
                cell.font = content_font
                cell.border = border
                cell.alignment = Alignment(horizontal='left', vertical='top', wrap_text=True, indent=1)
                current_row += 1
        
        return current_row

    def _add_excel_blockquote(self, ws, element, current_row, content_font, border):
        """添加 Excel 引用块"""
        text = element.get_text().strip()
        if text:
            cell = ws.cell(row=current_row, column=1, value=f"> {text}")
            cell.font = Font(size=10, italic=True)
            cell.border = border
            cell.alignment = Alignment(horizontal='left', vertical='top', wrap_text=True)
            # 设置背景色
            cell.fill = PatternFill(start_color="F0F0F0", end_color="F0F0F0", fill_type="solid")
            return current_row + 1
        return current_row

    def _add_excel_code_block(self, ws, element, current_row, content_font, border):
        """添加 Excel 代码块"""
        code = element.find('code')
        if code:
            text = code.get_text().strip()
        else:
            text = element.get_text().strip()
        
        if text:
            cell = ws.cell(row=current_row, column=1, value=text)
            cell.font = Font(name='Courier New', size=9)
            cell.border = border
            cell.alignment = Alignment(horizontal='left', vertical='top', wrap_text=True)
            # 设置背景色
            cell.fill = PatternFill(start_color="F8F9FA", end_color="F8F9FA", fill_type="solid")
            return current_row + 1
        return current_row

    def _add_excel_horizontal_rule(self, ws, current_row, border):
        """添加 Excel 分割线"""
        cell = ws.cell(row=current_row, column=1, value="─" * 50)
        cell.font = Font(size=10)
        cell.border = border
        cell.alignment = Alignment(horizontal='center')
        return current_row + 1

    def _add_excel_table(self, ws, element, current_row, header_font, header_fill, content_font, border):
        """添加 Excel 表格"""
        rows = element.find_all('tr')
        if rows:
            # 计算表格大小
            max_cols = max(len(row.find_all(['td', 'th'])) for row in rows)
            
            for i, row in enumerate(rows):
                cells = row.find_all(['td', 'th'])
                for j, cell in enumerate(cells):
                    if j < max_cols:
                        excel_cell = ws.cell(row=current_row + i, column=j + 1, value=cell.get_text().strip())
                        
                        # 设置样式
                        if cell.name == 'th':  # 表头
                            excel_cell.font = header_font
                            excel_cell.fill = header_fill
                        else:  # 数据单元格
                            excel_cell.font = content_font
                        
                        excel_cell.border = border
                        excel_cell.alignment = Alignment(horizontal='left', vertical='center', wrap_text=True)
            
            return current_row + len(rows)
        return current_row
    
    async def markdown_to_pdf(
        self,
        input_path: str,
        output_path: str,
        options: Optional[Dict[str, Any]] = None
    ) -> Dict[str, Any]:
        """
        将 Markdown 转换为 PDF
        
        Args:
            input_path: 输入 Markdown 文件路径
            output_path: 输出 PDF 文件路径
            options: 转换选项
            
        Returns:
            转换结果字典
        """
        try:
            if not DEPENDENCIES_AVAILABLE:
                return {
                    'success': False,
                    'error': 'Markdown 转换器依赖未安装，请安装 weasyprint 和 markdown'
                }
            
            logger.info(f"开始 Markdown 转 PDF: {input_path}")
            
            # 读取 Markdown 文件
            with open(input_path, 'r', encoding='utf-8') as f:
                md_content = f.read()
            
            # 转换 Markdown 为 HTML
            html_content = markdown.markdown(
                md_content,
                extensions=['tables', 'fenced_code', 'codehilite', 'toc']
            )
            
            # 构建完整的 HTML 文档
            css_style = options.get('css_style', self.default_css) if options else self.default_css
            
            full_html = f"""
            <!DOCTYPE html>
            <html>
            <head>
                <meta charset="utf-8">
                <title>Markdown Document</title>
                <style>
                    {css_style}
                </style>
            </head>
            <body>
                {html_content}
            </body>
            </html>
            """
            
            # 生成 PDF
            HTML(string=full_html).write_pdf(output_path)
            
            logger.info(f"Markdown 转 PDF 成功: {output_path}")
            return {
                'success': True,
                'output_path': output_path,
                'output_filename': Path(output_path).name
            }
            
        except Exception as e:
            logger.error(f"Markdown 转 PDF 失败: {e}")
            return {
                'success': False,
                'error': f'Markdown 转 PDF 失败: {str(e)}'
            }
    
    async def pdf_to_markdown(
        self,
        input_path: str,
        output_path: str,
        options: Optional[Dict[str, Any]] = None
    ) -> Dict[str, Any]:
        """
        将 PDF 转换为 Markdown
        
        Args:
            input_path: 输入 PDF 文件路径
            output_path: 输出 Markdown 文件路径
            options: 转换选项
            
        Returns:
            转换结果字典
        """
        try:
            logger.info(f"开始 PDF 转 Markdown: {input_path}")
            
            # 打开 PDF 文件
            doc = fitz.open(input_path)
            markdown_content = []
            
            # 处理每一页
            for page_num in range(len(doc)):
                page = doc.load_page(page_num)
                
                # 提取文本
                text = page.get_text()
                
                # 格式化文本为 Markdown
                formatted_text = self._format_text_to_markdown(text, page_num + 1)
                markdown_content.append(formatted_text)
            
            # 合并所有页面的内容
            final_markdown = '\n\n---\n\n'.join(markdown_content)
            
            # 保存 Markdown 文件
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(final_markdown)
            
            doc.close()
            
            logger.info(f"PDF 转 Markdown 成功: {output_path}")
            return {
                'success': True,
                'output_path': output_path,
                'output_filename': Path(output_path).name
            }
            
        except Exception as e:
            logger.error(f"PDF 转 Markdown 失败: {e}")
            return {
                'success': False,
                'error': f'PDF 转 Markdown 失败: {str(e)}'
            }
    
    def _format_text_to_markdown(self, text: str, page_num: int) -> str:
        """
        将文本格式化为 Markdown
        
        Args:
            text: 原始文本
            page_num: 页码
            
        Returns:
            格式化的 Markdown 文本
        """
        # 添加页面标题
        markdown_text = f"# 第 {page_num} 页\n\n"
        
        # 分割文本为段落
        paragraphs = text.split('\n\n')
        
        for paragraph in paragraphs:
            paragraph = paragraph.strip()
            if not paragraph:
                continue
            
            # 检测标题
            if self._is_heading(paragraph):
                markdown_text += f"## {paragraph}\n\n"
            # 检测列表项
            elif self._is_list_item(paragraph):
                markdown_text += f"- {paragraph}\n"
            # 检测代码块
            elif self._is_code_block(paragraph):
                markdown_text += f"```\n{paragraph}\n```\n\n"
            # 检测引用
            elif self._is_quote(paragraph):
                markdown_text += f"> {paragraph}\n\n"
            # 普通段落
            else:
                markdown_text += f"{paragraph}\n\n"
        
        return markdown_text
    
    def _is_heading(self, text: str) -> bool:
        """检测是否为标题"""
        # 简单的标题检测逻辑
        return len(text) < 100 and text.isupper() and not text.endswith('.')
    
    def _is_list_item(self, text: str) -> bool:
        """检测是否为列表项"""
        return text.startswith(('•', '-', '*', '1.', '2.', '3.'))
    
    def _is_code_block(self, text: str) -> bool:
        """检测是否为代码块"""
        # 简单的代码块检测
        return 'function' in text.lower() or 'class' in text.lower() or 'import' in text.lower()
    
    def _is_quote(self, text: str) -> bool:
        """检测是否为引用"""
        return text.startswith('"') or text.startswith('"')
    
    def get_supported_formats(self) -> Dict[str, list]:
        """
        获取支持的转换格式
        
        Returns:
            支持的转换格式字典
        """
        return {
            'input_formats': ['md', 'markdown', 'pdf'],
            'output_formats': ['pdf', 'md', 'markdown']
        }
    
    def is_supported_conversion(self, source_format: str, target_format: str) -> bool:
        """
        检查是否支持指定的转换
        
        Args:
            source_format: 源格式
            target_format: 目标格式
            
        Returns:
            是否支持
        """
        supported_formats = self.get_supported_formats()
        
        return (
            source_format in supported_formats['input_formats'] and
            target_format in supported_formats['output_formats']
        ) 